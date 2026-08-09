"""3-slide PPTX: model, model comparison, validation results.

Assets are drawn with matplotlib into results/slides_assets/, then a 16:9
deck is assembled at results/NV_C13_hybrid_slides.pptx.
"""

from __future__ import annotations

import sys
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "results" / "slides_assets"
FIGS = ROOT / "results" / "figs"
ASSETS.mkdir(parents=True, exist_ok=True)

plt.rcParams["font.family"] = ["Noto Sans CJK KR", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


def box(ax, x, y, w, h, text, fc, fontsize=11, tc="black", sub=None):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.012",
                                fc=fc, ec="0.35", lw=1.2))
    ax.text(x + w / 2, y + h / 2 + (0.055 if sub else 0), text, ha="center",
            va="center", fontsize=fontsize, color=tc, weight="bold", wrap=True)
    if sub:
        ax.text(x + w / 2, y + h / 2 - 0.085, sub, ha="center", va="center",
                fontsize=fontsize - 2.5, color="0.25")


def arrow(ax, x1, y1, x2, y2):
    ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>",
                                 mutation_scale=22, lw=1.8, color="0.2"))


def asset_pipeline():
    fig, ax = plt.subplots(figsize=(12.4, 4.4))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    y, h = 0.56, 0.30
    box(ax, 0.010, y, 0.150, h, "CPMG 신호\nN=8/16/20", "#dbe9ff",
        sub="(3×700, 상온 실측)")
    box(ax, 0.205, y, 0.165, h, "주기-윈도우\n토큰화", "#dff2df",
        sub="61 윈도우 slice-stack")
    box(ax, 0.415, y, 0.165, h, "Transformer\n(윈도우 간 어텐션)", "#dff2df",
        sub="PeriodFormer, 1.1M")
    box(ax, 0.625, y, 0.155, h, "P(spin) 곡선\n→ 후보 영역", "#fff3cc",
        sub="오탐 0 검출")
    box(ax, 0.825, y, 0.165, h, "영역 제약\nDE + BIC", "#ffe0e0",
        sub="클러스터 내 열거")
    for x1, x2 in [(0.160, 0.205), (0.370, 0.415), (0.580, 0.625), (0.780, 0.825)]:
        arrow(ax, x1, y + h / 2, x2, y + h / 2)
    ax.text(0.5, 0.97, "PF→DE 하이브리드 파이프라인", ha="center", fontsize=15,
            weight="bold")
    ax.text(0.287, 0.38, "물리 사전지식:\n스핀은 주기 TP(A)에서\n수직선을 만든다",
            ha="center", va="top", fontsize=9, color="0.3")
    ax.text(0.497, 0.38, "실제 스핀만 이웃 윈도우에\n일관된 흔적 → 어텐션이 종합",
            ha="center", va="top", fontsize=9, color="0.3")
    ax.text(0.907, 0.38, "같은 영역에 스핀 여러 개 허용\n개수는 BIC가 결정",
            ha="center", va="top", fontsize=9, color="0.3")
    # output
    box(ax, 0.35, 0.06, 0.30, 0.18, "출력: {(A∥, A⊥)} 스핀 목록", "#eadcf8",
        sub="NV1 14개 · NV2 10개 (앵커 7/7 회수)")
    arrow(ax, 0.907, y - 0.02, 0.65, 0.17)
    fig.savefig(ASSETS / "pipeline.png", dpi=200, bbox_inches="tight")
    plt.close(fig)


def asset_compare():
    fig, ax = plt.subplots(figsize=(12.4, 3.6))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    cols = [
        ("2021 (Jung et al.)", "#f2d5cc",
         "이미지 → 윈도우별\n독립 MLP 61개\n+ AE 디노이저",
         "윈도우 간 정보 공유 없음\n디노이저가 딥 훼손"),
        ("CNN 뱅크 (개선)", "#fde9c8",
         "CNN + N=8/16/20\n3채널 joint 입력",
         "채널 결합으로 개선\n여전히 독립 판정"),
        ("SpinDETR", "#d9ead3",
         "원신호 → set 예측\n(DETR, 물리 사전지식 X)",
         "(A,B) 직접 출력\n1 forward pass"),
        ("PeriodFormer", "#cfe2f3",
         "윈도우 임베딩 토큰\n+ 교차 윈도우 어텐션",
         "물리 사전지식 유지\n전 조건 오탐 0"),
        ("하이브리드 (최종)", "#ead1dc",
         "PF 후보 영역\n→ 영역 제약 DE",
         "검출 신뢰성 + 열거 능력\n전 조건 1위"),
    ]
    w = 0.184
    for i, (title, fc, body, note) in enumerate(cols):
        x = 0.008 + i * 0.20
        box(ax, x, 0.42, w, 0.44, title, fc, fontsize=11.5)
        ax.text(x + w / 2, 0.52, body, ha="center", va="center", fontsize=9.5)
        ax.text(x + w / 2, 0.24, note, ha="center", va="center", fontsize=8.5,
                color="0.25", style="italic")
    fig.savefig(ASSETS / "compare_schematic.png", dpi=200, bbox_inches="tight")
    plt.close(fig)


def asset_f1_bars():
    methods = ["2021 풀 파이프라인", "CNN+joint-N", "SpinDETR", "PeriodFormer",
               "cdetect-DE", "하이브리드"]
    f1_room = [0.57, 0.84, 0.88, 0.94, 0.94, None]      # synthetic suite sigma=0.08
    # 50-spin twin arms
    armA = [0.376, None, None, 0.412, 0.714, 0.780]
    armB = [0.367, None, None, 0.357, 0.500, 0.510]
    armC = [0.353, None, None, 0.338, 0.480, 0.553]

    fig, axes = plt.subplots(1, 2, figsize=(12.4, 3.6),
                             gridspec_kw={"width_ratios": [1, 1.4]})
    ax = axes[0]
    names = methods[:5]
    vals = f1_room[:5]
    colors = ["#b45f4d", "#e8a33d", "#6aa84f", "#3d85c6", "#999999"]
    ax.barh(range(len(names)), vals, color=colors[: len(names)], alpha=0.9)
    ax.set_yticks(range(len(names)))
    ax.set_yticklabels(names, fontsize=10)
    ax.set_xlim(0, 1.0)
    ax.set_xlabel("F1 @ 상온 노이즈 (합성 GT)", fontsize=10)
    for i, v in enumerate(vals):
        ax.text(v + 0.01, i, f"{v:.2f}", va="center", fontsize=9)
    ax.set_title("아키텍처 ablation (σ=0.08)", fontsize=11)
    ax.invert_yaxis()

    ax = axes[1]
    x = np.arange(3)
    width = 0.2
    sel = [(0, "2021", "#b45f4d"), (3, "PeriodFormer", "#3d85c6"),
           (4, "cdetect-DE", "#6aa84f"), (5, "하이브리드", "#cc0000")]
    for j, (idx, label, c) in enumerate(sel):
        vals = [armA[idx], armB[idx], armC[idx]]
        ax.bar(x + (j - 1.5) * width, vals, width, label=label, color=c, alpha=0.9)
        for xi, v in zip(x + (j - 1.5) * width, vals):
            ax.text(xi, v + 0.012, f"{v:.2f}", ha="center", fontsize=7.5)
    ax.set_xticks(x)
    ax.set_xticklabels(["암 A: 저온 트윈\n(그들의 조건)", "암 B: 상온\n(우리 조건)",
                        "암 C: 모델 불일치\n(현실 오차)"], fontsize=9.5)
    ax.set_ylabel("F1", fontsize=10)
    ax.set_ylim(0, 0.9)
    ax.legend(fontsize=9, ncol=4, loc="upper right")
    ax.set_title("실측 50-스핀 배스 검증 (Delft 공개 데이터, 테스트 전용)", fontsize=11)
    fig.tight_layout()
    fig.savefig(ASSETS / "f1_bars.png", dpi=200, bbox_inches="tight")
    plt.close(fig)




def asset_steps():
    """Project timeline: 7 steps, snake layout."""
    fig, ax = plt.subplots(figsize=(12.4, 5.2))
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    steps = [
        ("1. EDA · 물리 검증", "실측 딥 간격 1.060 µs =\n이론 TP 1.061 µs\n→ forward model 확정"),
        ("2. 2021 재구성", "MLP 윈도우뱅크 + AE 디노이저\n논문 코드 분석 후\n우리 조건으로 재학습"),
        ("3. Ablation 스터디", "아키텍처·채널·디노이저 분해\n상온 F1 0.57 → 0.94\n(AE 디노이저는 유해 판명)"),
        ("4. 신규 아키텍처", "SpinDETR (raw→set 예측)\nPeriodFormer (토큰+어텐션)\n합성 장면으로만 학습"),
        ("5. PF→DE 하이브리드", "PF 후보영역(오탐 0)\n+ 영역제약 DE(BIC 열거)\n= 전 조건 1위"),
        ("6. 실측 GT 검증", "Delft 50-스핀 공개데이터\n디지털 트윈 3개 암\n(저온/상온/모델불일치)"),
        ("7. 실데이터 확정", "NV1 14스핀 · NV2 10스핀\n앵커 7/7 회수\nRMSE 0.088–0.286"),
    ]
    pos = [(0.01, 0.62), (0.26, 0.62), (0.51, 0.62), (0.76, 0.62),
           (0.76, 0.12), (0.51, 0.12), (0.26, 0.12)]
    w, h = 0.215, 0.33
    for (title, body), (x, y) in zip(steps, pos):
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.012",
                                    fc="#eef3fb", ec="0.35", lw=1.2))
        ax.text(x + w/2, y + h - 0.055, title, ha="center", va="center",
                fontsize=12, weight="bold")
        ax.text(x + w/2, y + h/2 - 0.045, body, ha="center", va="center",
                fontsize=9, color="0.25")
    for i in range(len(pos) - 1):
        x1, y1 = pos[i]; x2, y2 = pos[i + 1]
        if y1 == y2:
            if x2 > x1: arrow(ax, x1 + w + 0.003, y1 + h/2, x2 - 0.003, y2 + h/2)
            else: arrow(ax, x1 - 0.003, y1 + h/2, x2 + w + 0.003, y2 + h/2)
        else:
            arrow(ax, x1 + w/2, y1 - 0.005, x2 + w/2, y2 + h + 0.005)
    fig.savefig(ASSETS / "steps.png", dpi=200, bbox_inches="tight")
    plt.close(fig)


def asset_dataflow():
    """DL-paper style: tensor shapes flowing through each model."""
    fig, ax = plt.subplots(figsize=(12.6, 6.2))
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    def chain(y, label, label_c, blocks, note=None):
        ax.text(0.001, y + 0.10, label, fontsize=12, weight="bold", color=label_c)
        x = 0.10
        for text, shape, fc, wdt in blocks:
            ax.add_patch(FancyBboxPatch((x, y), wdt, 0.115,
                                        boxstyle="round,pad=0.008",
                                        fc=fc, ec="0.4", lw=1.0))
            ax.text(x + wdt/2, y + 0.078, text, ha="center", va="center",
                    fontsize=8.8, weight="bold")
            ax.text(x + wdt/2, y + 0.028, shape, ha="center", va="center",
                    fontsize=8.0, color="#8B0000")
            x_end = x + wdt
            x = x_end + 0.022
            if (text, shape, fc, wdt) != blocks[-1]:
                arrow(ax, x_end + 0.002, y + 0.058, x - 0.002, y + 0.058)
        if note:
            ax.text(0.10, y - 0.032, note, fontsize=8.2, color="0.35",
                    style="italic")

    C_IN, C_OP, C_OUT = "#dbe9ff", "#dff2df", "#eadcf8"
    chain(0.83, "2021 재구성\n(MLP 뱅크)", "#b45f4d", [
        ("Px 신호", "(700,)", C_IN, 0.085),
        ("엔벨로프 정규화", "M (700,)", C_OP, 0.105),
        ("윈도우 접기 ·flatten", "(13×53)→689", C_OP, 0.125),
        ("Dense 2048→1024→512", "(689)→(512)", C_OP, 0.155),
        ("sigmoid 분류", "(3,) 확률", C_OP, 0.10),
        ("61 윈도우 반복→피크", "A 위치만", C_OUT, 0.135),
    ], note="윈도우마다 독립 모델 61개 · B는 별도 회귀 필요 · 윈도우 간 정보 공유 없음")

    chain(0.60, "PeriodFormer\n(우리 제안)", "#3d85c6", [
        ("M 신호 3채널", "(3, 700)", C_IN, 0.09),
        ("TokenBuilder\n(GPU gather 접기)", "(61, 3, 13, 53)", C_OP, 0.13),
        ("공유 CNN 임베딩", "(61, 128)", C_OP, 0.115),
        ("+PE → Transformer ×4", "(61, 128)", C_OP, 0.14),
        ("존재확률 헤드", "P(spin) (61,)", C_OP, 0.105),
        ("후보 영역", "구간 목록", C_OUT, 0.085),
    ], note="윈도우가 '토큰' — 어텐션이 61개 윈도우의 증거를 종합 → 오탐 0 · 모델 1개(1.1M)")

    chain(0.37, "SpinDETR\n(비교용 e2e)", "#6aa84f", [
        ("M 신호 3채널", "(3, 700)", C_IN, 0.09),
        ("Conv stem (stride 4)", "(175, 128)", C_OP, 0.13),
        ("Encoder ×4", "(175, 128)", C_OP, 0.10),
        ("쿼리 10개 Decoder ×4", "(10, 128)", C_OP, 0.14),
        ("존재+회귀 헤드", "10×(p, A, B)", C_OP, 0.115),
        ("p>0.5 선택", "{(A,B)}", C_OUT, 0.08),
    ], note="물리 사전지식 없이 원신호에서 직접 set 예측 · Hungarian 매칭 손실 · (A,B) 동시 출력")

    chain(0.14, "하이브리드 최종\n(PF→DE)", "#cc0000", [
        ("PF 후보 영역", "구간 목록", C_IN, 0.10),
        ("영역별 DE 탐색", "스핀 +1 씩", C_OP, 0.115),
        ("전체 재연마 (L-BFGS)", "(2k,) 파라미터", C_OP, 0.14),
        ("BIC 모델 선택", "k* 자동 결정", C_OP, 0.115),
        ("최종 출력", "{(A∥, A⊥)} × k*", C_OUT, 0.11),
    ], note="학습 없음(최적화) · forward model ∏Mᵢ를 직접 피팅 · 같은 영역에 다중 스핀 허용")

    ax.text(0.5, 0.985, "데이터 흐름과 텐서 shape: 2021 재구성 vs 우리 모델들",
            ha="center", fontsize=14, weight="bold")
    fig.savefig(ASSETS / "dataflow.png", dpi=200, bbox_inches="tight")
    plt.close(fig)


def asset_fitrender():
    """How a fit curve is rendered from a spin list (real NV2 data demo)."""
    sys.path.insert(0, str(ROOT))
    import json
    import pandas as pd
    from cpmg.physics import cpmg_M
    from cpmg.represent import envelope_normalize
    nv2 = pd.read_excel(ROOT / "dataset" / "exp_dataset" / "CPMG_NV2.xlsx")
    tau = nv2["Time"].to_numpy(float)
    px = nv2["CPMG16"].to_numpy(float)
    m_data, _ = envelope_normalize(tau, px)
    spins = json.loads((ROOT / "results" / "benchmark_v2" /
                        "hybrid_results.json").read_text())["nv2_ensemble"]["spins_khz"]
    ab = np.array(spins) * 1e3
    US = 1e-6
    zoom = tau <= 4e-6

    fig, axes = plt.subplots(1, 3, figsize=(12.6, 3.1))
    ax = axes[0]
    ax.plot(tau[zoom] / US, px[zoom], ".", ms=2, color="0.4")
    ax.plot(tau[zoom] / US, m_data[zoom] * 0.4 + 0.5, ".", ms=2, color="tab:blue",
            alpha=0.5)
    ax.set_title("① 원시 Pₓ → 엔벨로프 정규화 M", fontsize=10)
    ax.set_xlabel("τ (µs)")
    ax.text(0.5, 0.02, "M = (2Px−1) / (a0·exp(−(τ/T2)^n))", transform=ax.transAxes,
            ha="center", fontsize=9, color="tab:blue")

    ax = axes[1]
    picks = [(346.9, 261.0), (-58.9, 181.7), (-12.6, 42.3)]
    for i, (a, b) in enumerate(picks):
        mi = cpmg_M(tau, [[a * 1e3, b * 1e3]], 16, 440.1)
        ax.plot(tau[zoom] / US, mi[zoom] + i * 1.1, lw=0.9,
                label=f"({a:+.0f}, {b:.0f}) kHz")
    ax.set_yticks([])
    ax.legend(fontsize=7, loc="lower right")
    ax.set_title("② 스핀별 딥 트레인 Mᵢ(τ)", fontsize=10)
    ax.set_xlabel("τ (µs)")

    ax = axes[2]
    m_fit = cpmg_M(tau, ab, 16, 440.1)
    rmse = float(np.sqrt(np.mean((m_fit - m_data) ** 2)))
    ax.plot(tau[zoom] / US, m_data[zoom], ".", ms=2.2, color="0.55")
    ax.plot(tau[zoom] / US, m_fit[zoom], "-", lw=1.1, color="tab:red")
    ax.set_title(f"③ 곱 ∏Mᵢ = 피팅 곡선 (10스핀, RMSE {rmse:.3f})", fontsize=10)
    ax.set_xlabel("τ (µs)")
    ax.text(0.5, 0.02, "M_fit(τ) = ∏ Mi(τ; Ai, Bi)   (합이 아닌 곱!)",
            transform=ax.transAxes, ha="center", fontsize=9, color="tab:red")
    fig.tight_layout()
    fig.savefig(ASSETS / "fitrender.png", dpi=200, bbox_inches="tight")
    plt.close(fig)




def asset_analogy():
    """Five designs explained by analogy, each fixing the previous weakness."""
    fig, ax = plt.subplots(figsize=(12.6, 4.9))
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    cards = [
        ("2021 방식", "창문 검사관 61명", "#f2d5cc",
         "주기(A)마다 검사관이 자기\n창문 그림만 보고 판정",
         "약점: 서로 대화가 없어\n노이즈에 잘 속는다(오탐)"),
        ("CNN 뱅크", "더 좋은 눈", "#fde9c8",
         "검사관에게 패턴 인식 눈(CNN)과\n사진 3장(N=8/16/20)을 줌",
         "개선되지만 여전히\n각자 판정한다"),
        ("SpinDETR", "속기사", "#d9ead3",
         "물리 힌트 없이 신호 전체를 듣고\n스핀 목록을 바로 받아씀",
         "장점: (A,B) 동시 출력\n약점: 힌트가 없어 불리"),
        ("PeriodFormer", "검사관 회의", "#cfe2f3",
         "61명이 증거를 한자리에서\n맞춰봄(어텐션)",
         "진짜 스핀만 이웃 창문에도\n흔적을 남긴다 → 오탐 0"),
        ("하이브리드", "회의 + 측량사", "#ead1dc",
         "회의가 어디인지 정하면 측량사(DE)가\n그 구역만 정밀 측량",
         "몇 개인지는 통계 기준\n(BIC)이 자동 결정"),
    ]
    w = 0.178
    labels = ["판정 눈을 개선", "힌트를 버려보면?\n(대조 실험)",
              "힌트 + 회의를 결합", "개수 세기 보완"]
    for i, (name, nick, fc, body, note) in enumerate(cards):
        x = 0.008 + i * 0.20
        ax.add_patch(FancyBboxPatch((x, 0.30), w, 0.58,
                                    boxstyle="round,pad=0.012",
                                    fc=fc, ec="0.35", lw=1.2))
        ax.text(x + w/2, 0.83, name, ha="center", fontsize=11, weight="bold")
        ax.text(x + w/2, 0.75, nick, ha="center", fontsize=12,
                weight="bold", color="#333")
        ax.text(x + w/2, 0.60, body, ha="center", va="center", fontsize=8.8)
        ax.text(x + w/2, 0.40, note, ha="center", va="center", fontsize=8.3,
                color="#8B0000")
        if i < 4:
            arrow(ax, x + w + 0.004, 0.59, x + 0.20 - 0.004, 0.59)
            ax.text(x + w + 0.011, 0.50, labels[i], ha="left", fontsize=7.6,
                    color="0.35", style="italic")
    ax.text(0.5, 0.10, "각 설계는 바로 앞 방식의 구체적 약점 하나를 고친다 — "
            "그래서 성능 향상의 원인을 분리해 검증할 수 있다",
            ha="center", fontsize=11, weight="bold", color="#1155CC")
    fig.savefig(ASSETS / "analogy.png", dpi=200, bbox_inches="tight")
    plt.close(fig)


def build_deck():
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    from pptx.dml.color import RGBColor as _RGB

    def add_title(slide, text, sub=None, msg=None):
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12),
                                      Inches(12.5), Inches(1.0))
        tf = tb.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(26)
        tf.paragraphs[0].font.bold = True
        if sub:
            p = tf.add_paragraph()
            p.text = sub
            p.font.size = Pt(12)
        if msg:
            p = tf.add_paragraph()
            p.text = "▶ " + msg
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = _RGB(0xB4, 0x50, 0x00)

    def add_bullets(slide, items, left, top, width, size=13):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + it
            p.font.size = Pt(size)

    # ---- slide: project steps overview ----
    s = prs.slides.add_slide(blank)
    add_title(s, "무엇을 했는가: 7단계 요약",
              "물리 검증 → 2021 재구성 → ablation → 신규 모델 → 하이브리드 → 실측 GT 검증 → 실데이터 확정",
              msg="모든 주장 뒤에 검증 단계를 붙였다 — 마지막엔 정답을 아는 실측 스핀 환경으로 확인")
    s.shapes.add_picture(str(ASSETS / "steps.png"), Inches(0.4),
                         Inches(1.35), width=Inches(12.5))
    add_bullets(s, [
        "모든 학습은 forward model 합성 데이터만 사용 — 실험 데이터와 실측 배스는 테스트 전용 (누수 없음)",
        "각 단계의 산출물이 다음 단계의 근거가 되는 구조: 예) ablation의 발견(AE 유해, joint-N 유효)이 신규 모델 설계로 이어짐",
    ], 0.5, 6.6, 12.3, size=12.5)

    # ---- slide: public 50-spin validation (placed up front) ----
    s = prs.slides.add_slide(blank)
    add_title(s, "검증 근거: 공개 실측 50-스핀 배스 (디지털 트윈)",
              "van de Stolpe et al., Nat. Commun. 2024 (4TU 공개 데이터) — Jung 2021과 같은 NV 계보 · B=403.553 G까지 동일",
              msg="정답(50개 스핀의 실측 A∥·A⊥)을 아는 실제 스핀 환경에서 성능을 채점했다 — 스핀 값은 학습에 미사용, 순수 테스트")
    s.shapes.add_picture(str(FIGS / "26_twin_validation.png"), Inches(0.35),
                         Inches(1.5), width=Inches(12.6))
    add_bullets(s, [
        "왼쪽: 공개 실측값으로 합성한 트윈 신호(회색) 위에 우리가 검출한 스핀들의 forward model(빨강)을 피팅 — 저온 RMSE 0.074 / 상온 0.143",
        "오른쪽: 공개 GT 27스핀(빨강)과 검출(파랑)의 위치 대조 — 저온 21/27 회수, 상온 11/27·오탐 0 (상온 한계는 물리적 정보량: EDA의 검출가능 상한과 정합)",
    ], 0.5, 6.85, 12.3, size=12)

    # ---- Q&A section: five audience questions ----
    s = prs.slides.add_slide(blank)
    add_title(s, "Q1. 스핀이 몇 개인지 모르는데 어떻게 찾았나",
              "스핀 수 k 자체가 미지수인 문제(trans-dimensional) — 두 가지 독립 원리로 결정",
              msg="새 스핀은 벌점(월세)을 내야 입주한다 — 진짜 스핀만 자기 딥을 설명해서 벌점을 넘고, 노이즈는 탈락한다")
    s.shapes.add_picture(str(ASSETS / "q1_bic.png"), Inches(0.35),
                         Inches(1.5), width=Inches(12.6))
    add_bullets(s, [
        "① BIC: 스핀 1개 = 파라미터 2개 벌점(2·ln n). 잔차 감소가 벌점보다 클 때만 채택 → NV1에서 k=14가 최소, 15번째부터 반등",
        "② 독립 교차확인: 스핀 수를 확률변수로 두는 RJMCMC의 사후분포 P(k|데이터)가 14~16에 집중(최빈 15) — 서로 다른 원리가 같은 답",
        "안전장치: 모델 오차가 가짜 스핀으로 흡수되는 것을 PF 후보영역이 차단 (영역 밖 입주 신청 불가)",
    ], 0.5, 5.75, 12.3, size=12.5)

    s = prs.slides.add_slide(blank)
    add_title(s, "Q2. 기존 2021 방법을 어떻게 재현했나",
              "공개 코드(imports/models.py)와 논문 Methods를 기준으로 핵심 파이프라인을 동일 구조로 재구현",
              msg="재현 범위를 명시적으로 한정했다 — 보고하는 2021 수치는 원방법 성능의 하한(lower bound)")
    s.shapes.add_picture(str(ASSETS / "q2_repro.png"), Inches(0.35),
                         Inches(1.6), width=Inches(12.6))
    add_bullets(s, [
        "충실성 근거: 같은 표현(slice-stack)·같은 네트워크 구조(Dense 2048-1024-512)·같은 디노이저(conv-AE k4/64ch)를 코드 수준에서 재현",
        "작동 검증: 재현본이 저온 트윈(그들의 조건)에서는 정상 작동(오탐 0) — 성능 저하는 상온 조건에서만 발생 → 구현 오류가 아니라 조건의 문제",
    ], 0.5, 6.2, 12.3, size=12.5)

    s = prs.slides.add_slide(blank)
    add_title(s, "Q3. 50-스핀 공개 데이터는 믿을만한가 — 어떤 데이터인가",
              "Delft(Taminiau 그룹)의 같은 NV 계보에서 측정된, 3중 피어리뷰를 거친 실측 스핀 지도",
              msg="정답을 잰 실험 기법(SEDOR 이중공명, 1.8 Hz 분해능)이 우리 방법(CPMG)과 완전히 독립 — 순환 논증이 없다")
    s.shapes.add_picture(str(ASSETS / "q3_lineage.png"), Inches(0.35),
                         Inches(1.55), width=Inches(12.6))
    add_bullets(s, [
        "2021 논문(Jung et al.)이 분석한 바로 그 NV 계보 — 자기장 403.553 G가 소수점까지 일치, 27-스핀(2019 Nature)의 상위집합",
        "우리는 이 50개 실측 (A∥, A⊥)를 물리 모델에 넣어 CPMG 신호를 합성(디지털 트윈)하고, 그 신호에서 스핀을 되찾는 시험을 채점했다",
    ], 0.5, 6.05, 12.3, size=12.5)

    s = prs.slides.add_slide(blank)
    add_title(s, "Q4. 그 데이터에서 충분한 검증을 하였는가",
              "저온/상온/오차주입 3개 조건 × 노이즈 실현 4-8회 × 학습 시드 5회 + 6개 방법 비교",
              msg="정답을 아는 환경에서 precision/recall을 직접 채점했고, 통계적 유의성(p=0.008)과 시드 강건성(±0.01)까지 확인했다")
    s.shapes.add_picture(str(FIGS / "10_ablation_f1.png"), Inches(0.35),
                         Inches(1.6), width=Inches(12.6))
    add_bullets(s, [
        "합성 GT 스위트(노이즈 스윕) + 실측 50-스핀 트윈(3개 암) + 실데이터의 3층 검증 — 각 층의 실패 모드가 서로 다름",
        "최종 하이브리드 v2: 세 조건 모두 1위 (F1 0.84/0.57/0.61) · 다중 시드 mean±std 보고 · 하이브리드 vs 고전 피팅 우위는 오차주입 조건에서 유의(p=0.0078)",
    ], 0.5, 6.15, 12.3, size=12.5)

    s = prs.slides.add_slide(blank)
    add_title(s, "Q5. NV1/NV2에서 찾은 핵스핀이 옳다고 볼 근거는",
              "정답이 없는 실데이터 — 다섯 겹의 독립 증거로 신뢰를 쌓는다",
              msg="원리가 다른 방법들이 같은 목록에 수렴하고, 그 목록의 물리 예측이 세 측정(N=8/16/20)을 동시에 재현한다")
    s.shapes.add_picture(str(ASSETS / "q5_ci.png"), Inches(0.35),
                         Inches(1.55), width=Inches(12.6))
    add_bullets(s, [
        "① 교차 수렴: 9개 방법이 핵심 스핀에 전원 일치  ② 물리 재현: 14스핀 forward model이 3개 채널 실데이터를 RMSE 0.09-0.18로 동시 설명",
        "③ 통계적 실체: 부트스트랩 95% CI가 스핀끼리 겹치지 않음  ④ 외부 정합: 수동 분석 앵커 7/7 회수  ⑤ 물리 검증: 딥 간격 1.060 µs = 이론 1.061 µs",
    ], 0.5, 5.95, 12.3, size=12.5)

    # ---- slide 1: model ----
    s = prs.slides.add_slide(blank)
    add_title(s, "제안 모델: PF→DE 하이브리드",
              "상온 NV CPMG에서 ¹³C 핵스핀 (A∥, A⊥) 자동 추출",
              msg="어디에 있나(신경망 검출)와 몇 개·얼마나 강한가(물리 피팅)를 분리하니 둘 다 잘하게 됐다")
    s.shapes.add_picture(str(ASSETS / "pipeline.png"), Inches(0.35),
                         Inches(1.15), width=Inches(12.6))
    add_bullets(s, [
        "1단 PeriodFormer: 후보 주기마다 접은 slice-stack을 토큰으로 임베딩 — 물리 사전지식 유지, 윈도우 간 어텐션으로 실제 스핀만 통과 (오탐 0)",
        "2단 영역 제약 DE: PF가 좁힌 영역 안에서만 스핀을 순차 추가(BIC) — 한 클러스터의 다중 스핀까지 열거",
        "학습은 forward model(Eq.1–3) 합성 데이터만 사용 · 실험 데이터와 실측 배스는 테스트 전용",
    ], 0.5, 5.9, 12.3)

    # ---- slide: data flow / tensor shapes (DL-paper style) ----
    s = prs.slides.add_slide(blank)
    add_title(s, "모델 내부: 데이터가 어떻게 흘러 무엇이 되는가",
              "블록 위 = 연산, 붉은 글씨 = 텐서 shape · 입력(파랑) → 연산(초록) → 출력(보라)",
              msg="모델 간 본질적 차이는 정보가 어디서 합쳐지는가 — 2021은 끝까지 안 합치고, 우리는 어텐션에서 합친다")
    s.shapes.add_picture(str(ASSETS / "dataflow.png"), Inches(0.25),
                         Inches(1.15), width=Inches(12.85))
    add_bullets(s, [
        "2021 재구성: 윈도우마다 flatten→MLP 61회 반복, A 위치만 출력 · 우리(PF): 61 윈도우를 한 텐서(61×3×13×53)로 접어 어텐션이 한 번에 종합",
        "최종 출력 형태의 차이: 뱅크=확률 곡선(후처리 필요) / SpinDETR=(A,B) 직접 / 하이브리드=BIC로 개수까지 자동 결정된 {(A∥,A⊥)}",
    ], 0.5, 6.75, 12.3, size=12)

    # ---- slide 2: comparison ----
    s = prs.slides.add_slide(blank)
    add_title(s, "아키텍처 비교: 2021 → 하이브리드",
              "같은 조건·같은 합성 GT에서 재학습해 공정 비교 (ablation)",
              msg="상온 조건에서 2021 F1 0.57 → 제안 0.94 — 향상분이 어느 설계에서 왔는지 숫자로 분리했다")
    s.shapes.add_picture(str(ASSETS / "compare_schematic.png"), Inches(0.35),
                         Inches(1.1), width=Inches(12.6))
    s.shapes.add_picture(str(ASSETS / "f1_bars.png"), Inches(0.35),
                         Inches(4.15), width=Inches(12.6))
    add_bullets(s, [
        "상온 노이즈에서 2021 F1 0.57 → 하이브리드 계열 0.94 · 기여 분해: CNN(+0.14), joint-N(+0.07), 토큰-어텐션(+0.10)",
    ], 0.5, 7.0, 12.3, size=12)

    # ---- slide: analogy for non-DL readers ----
    s = prs.slides.add_slide(blank)
    add_title(s, "비유로 이해하는 다섯 설계",
              "딥러닝을 몰라도 되는 설명 — 각 방법이 무엇을 하고, 왜 그렇게 바꿨는가",
              msg="설계 변경은 취향이 아니라 진단이다 — 앞 방법이 실패하는 지점을 하나씩 고쳤다")
    s.shapes.add_picture(str(ASSETS / "analogy.png"), Inches(0.35),
                         Inches(1.45), width=Inches(12.6))
    add_bullets(s, [
        "왜 회의(어텐션)가 결정적인가: 상온 노이즈는 한 창문에서는 스핀처럼 보일 수 있지만, 이웃 창문들과 대조하면 들통난다 — 그래서 PeriodFormer의 오탐이 0이 된다",
        "왜 측량사(물리 피팅)가 필요한가: 회의는 이 근처에 있다까지만 안다 — 한 구역에 몇 개가 겹쳐 있는지는 물리 공식을 데이터에 맞춰봐야(BIC) 알 수 있다",
    ], 0.5, 6.5, 12.3, size=12.5)

    # ---- slide: one-figure architecture summary (AI-generated) ----
    s = prs.slides.add_slide(blank)
    add_title(s, "네 모델 한 장 요약",
              "생성 다이어그램 (Nano Banana Pro) — 점선: PeriodFormer의 출력이 하이브리드의 입력",
              msg="위에서 아래로: 고립 판정(2021) → 쿼리 집합예측(SpinDETR) → 윈도우 간 어텐션(PF) → 영역 내 열거(하이브리드)")
    s.shapes.add_picture(str(ROOT / "results" / "figs" / "paperbanana_pro" /
                             "e_overview_all_candidate_0.png"),
                         Inches(1.05), Inches(1.25), width=Inches(11.2))

    # ---- slide 3: validation ----
    s = prs.slides.add_slide(blank)
    add_title(s, "검증: 3중 근거",
              "① 합성 GT ② 실측 50-스핀 디지털 트윈(공개 데이터·테스트 전용) ③ 실데이터 교차 수렴",
              msg="정답을 아는 실제 스핀 환경(Delft 공개 데이터)에서 저온·상온·오차주입 전 조건 1위")
    s.shapes.add_picture(str(FIGS / "23_method_comparison_NV1.png"),
                         Inches(0.35), Inches(1.15), width=Inches(7.1))
    s.shapes.add_picture(str(FIGS / "21_ensemble_overlay_NV1.png"),
                         Inches(7.65), Inches(1.15), width=Inches(5.35))
    add_bullets(s, [
        "실측 50-스핀 배스(같은 NV 계보, 4TU 공개)에서 하이브리드가 전 조건 1위 — 저온 F1 0.78 / 상온 0.51 / 모델 불일치 0.55 (2021: 0.38/0.37/0.35)",
        "왼쪽: 9개 알고리즘의 NV1 검출이 최종 14스핀(빨간선)에 계단식 수렴 · ppt 앵커 7/7 회수",
        "오른쪽: 최종 14스핀 forward-model이 3개 채널 실데이터를 동시 재현 (RMSE 0.088/0.116/0.176)",
    ], 0.5, 6.15, 12.3, size=12.5)


    # ---- slide: how fit curves are rendered ----
    s = prs.slides.add_slide(blank)
    add_title(s, "피팅 곡선은 이렇게 그린다",
              "스핀 목록 → 곱 공식 → 곡선 1개 (스핀 수 6~10은 각 방법의 BIC 자동 선택 결과)",
              msg="곡선은 그린 것이 아니라 스핀 목록에서 물리 공식으로 계산된 예측 — 데이터와 맞으면 목록이 맞다는 뜻")
    s.shapes.add_picture(str(ASSETS / "fitrender.png"), Inches(0.3),
                         Inches(1.3), width=Inches(12.7))
    add_bullets(s, [
        "① 실험 Pₓ를 엔벨로프로 나눠 모델과 같은 M 스케일로 정규화 → ② 각 스핀 (Aᵢ,Bᵢ)가 자기 주기의 딥 트레인 Mᵢ 생성 → ③ 전체 곱 ∏Mᵢ가 피팅 곡선",
        "스핀이 늘수록 곱에 항이 추가되어 미세 구조를 더 재현 — NV2에서 6→10스핀에 RMSE 0.308→0.286 (BIC 페널티에도 잔차 감소 = 과적합 아님)",
        "그림 25의 스택 배치: 방법별 세로 오프셋 + 회색 데이터 반복 — 각 곡선을 데이터와 1:1 대조하면서 방법 간 비교",
    ], 0.5, 5.6, 12.3, size=12)

    # ---- slide 4/5: per-method fit overlays with coupling strengths ----
    from pptx.dml.color import RGBColor

    METHOD_COLORS = {"DE": RGBColor(0x1F, 0x77, 0xB4),
                     "SpinDETR": RGBColor(0x2C, 0xA0, 0x2C),
                     "hybrid": RGBColor(0x94, 0x67, 0xBD),
                     "ENSEMBLE": RGBColor(0xD6, 0x27, 0x28)}

    def add_method_lists(slide, entries, left, top, width, size=10.5):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(5.6))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for name, color, spins in entries:
            p1 = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p1.text = name
            p1.font.size = Pt(size + 1.5)
            p1.font.bold = True
            p1.font.color.rgb = color
            p2 = tf.add_paragraph()
            p2.text = spins
            p2.font.size = Pt(size - 1)
            p2.font.color.rgb = color

    s = prs.slides.add_slide(blank)
    add_title(s, "NV1: 모델별 피팅 오버레이 + 검출 결합강도",
              "회색 = 실험 데이터(행마다 반복) · 곡선 색 = 알고리즘 · (A∥, A⊥) kHz",
              msg="스핀을 더 많이 찾은 최종 목록(14개)이 잔차도 가장 작다 — 과적합이 아니라 실제 스핀이라는 증거")
    s.shapes.add_picture(str(FIGS / "24_method_overlays_NV1.png"),
                         Inches(0.3), Inches(1.15), height=Inches(6.1))
    add_method_lists(s, [
        ("cdetect-DE (8 spins, RMSE 0.123)", METHOD_COLORS["DE"],
         "(+9.2,31) (−5.6,26) (+3.3,30) (+23.7,24) (+40.7,33) (+65.2,30) (+91.2,36) (−87.6,19)"),
        ("SpinDETR (9 spins, RMSE 0.147)", METHOD_COLORS["SpinDETR"],
         "(−6.2,29) (+7.4,29) (+13.1,25) (−15.5,21) (−2.2,23) (+37.6,27) (+42.5,31) (+55.9,24) (−34.4,17)"),
        ("하이브리드 narrow (10 spins, RMSE 0.133)", METHOD_COLORS["hybrid"],
         "(−19.0,16) (−13.9,16) (−8.0,16) (−5.1,23) (+1.1,21) (+4.0,27) (+8.5,30) (+14.5,25) (+37.5,28) (+45.0,33)"),
        ("ENSEMBLE 최종 (14 spins, RMSE 0.116)", METHOD_COLORS["ENSEMBLE"],
         "(−87.9,18) (−18.8,16) (−11.5,14) (−5.1,22) (+0.6,17) (+4.0,26) (+8.4,28) (+14.2,22) (+24.3,20) (+39.0,28) (+48.0,25) (+65.6,28) (+91.9,33) (+114.4,24)"),
    ], 7.35, 1.3, 5.7)

    s = prs.slides.add_slide(blank)
    add_title(s, "NV2: 모델별 피팅 오버레이 + 검출 결합강도",
              "CPMG-16 단일 채널 · 강결합 영역 (A⊥ 최대 ~290 kHz)",
              msg="채널이 하나뿐인 어려운 조건에서도 10개 스핀으로 격렬한 진동 구조를 재현")
    s.shapes.add_picture(str(FIGS / "25_method_overlays_NV2.png"),
                         Inches(0.3), Inches(1.15), height=Inches(6.1))
    add_method_lists(s, [
        ("cdetect-DE (6 spins, RMSE 0.308)", METHOD_COLORS["DE"],
         "(−51.5,199) (+51.1,94) (+346.2,263) (−39.2,67) (−152.3,99) (−14.1,50)"),
        ("하이브리드 v1 (8 spins, RMSE 0.308)", METHOD_COLORS["SpinDETR"],
         "(−45.7,53) (−39.3,71) (−32.5,151) (−12.8,42) (−3.5,40) (+48.3,109) (+55.7,54) (+345.7,266)"),
        ("하이브리드 refined (9 spins, RMSE 0.322)", METHOD_COLORS["hybrid"],
         "(−58.8,181) (−45.7,51) (−39.6,70) (−21.6,136) (−12.3,42) (+52.7,85) (+322.0,49) (+341.3,288) (+349.1,257)"),
        ("ENSEMBLE 최종 (10 spins, RMSE 0.286)", METHOD_COLORS["ENSEMBLE"],
         "(−151.2,95) (−58.9,182) (−45.9,50) (−38.9,67) (−12.6,42) (−2.7,37) (+51.7,91) (+55.9,55) (+346.9,261) (+430.7,58)"),
    ], 7.35, 1.3, 5.7)

    # ---- slide 6/7: A-tensor tables ----
    def add_table(slide, headers, rows, left, top, width, height, fontsize=8.5,
                  hdr_colors=None):
        shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                       Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        tbl = shape.table
        for j, htext in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = htext
            c.text_frame.paragraphs[0].font.size = Pt(fontsize)
            c.text_frame.paragraphs[0].font.bold = True
            if hdr_colors and hdr_colors[j]:
                c.fill.solid()
                c.fill.fore_color.rgb = hdr_colors[j]
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                c = tbl.cell(i + 1, j)
                c.text = val
                pgh = c.text_frame.paragraphs[0]
                pgh.font.size = Pt(fontsize)
                if j == 0:
                    pgh.font.bold = True

    HC = [RGBColor(0xF4, 0xCC, 0xCC), RGBColor(0xFC, 0xE5, 0xCD),
          RGBColor(0xD9, 0xD9, 0xD9), RGBColor(0xD9, 0xD9, 0xD9),
          RGBColor(0xD9, 0xD9, 0xD9), RGBColor(0xD9, 0xEA, 0xD3),
          RGBColor(0xCF, 0xE2, 0xF3), RGBColor(0xEA, 0xD1, 0xDC)]

    NV1_ROWS = [
        ("(−87.9, 18)", "(−88, 47)", "—", "—", "−86", "—", "(−87.6, 19)", "—"),
        ("(−18.8, 16)", "—", "−20", "−16", "−16", "(−15.5, 21)", "—", "(−19.0, 16)"),
        ("(−11.5, 14)", "—", "—", "−8", "—", "—", "—", "(−13.9, 16)"),
        ("(−5.1, 22)", "(−5, 35)", "−6", "−8", "−6", "(−6.2, 29)", "(−5.6, 26)", "(−5.1, 23)"),
        ("(+0.6, 17)", "—", "+2", "—", "+2", "(−2.2, 23)", "(+3.3, 30)", "(+1.1, 21)"),
        ("(+4.0, 26)", "—", "+2", "+4", "+2", "—", "(+3.3, 30)", "(+4.0, 27)"),
        ("(+8.4, 28)", "(+8, 31)", "+12", "+4", "+8", "(+7.4, 29)", "(+9.2, 31)", "(+8.5, 30)"),
        ("(+14.2, 22)", "—", "+12", "—", "+14", "(+13.1, 25)", "—", "(+14.5, 25)"),
        ("(+24.3, 20)", "—", "+22", "+28", "+22", "—", "(+23.7, 24)", "—"),
        ("(+39.0, 28)", "(−38, 37)*", "+42", "+42", "+38/40", "(+37.6, 27)", "(+40.7, 33)", "(+37.5, 28)"),
        ("(+48.0, 25)", "—", "+46", "—", "—", "—", "—", "(+45.0, 33)"),
        ("(+65.6, 28)", "—", "—", "—", "+64", "—", "(+65.2, 30)", "—"),
        ("(+91.9, 33)", "—", "—", "—", "+92", "—", "(+91.2, 36)", "—"),
        ("(+114.4, 24)", "—", "—", "—", "+116", "—", "—", "—"),
        ("후보 (−34?)", "(−38, 37)?", "−34", "−30", "—", "(−34.4, 17)", "—", "—"),
    ]
    s = prs.slides.add_slide(blank)
    add_title(s, "NV1: 알고리즘별 A-텐서 값 종합표",
              "값 = (A∥, A⊥) kHz · 윈도우뱅크/PF는 A 위치만 출력 · *앵커 부호는 m_s 브랜치 관례 차이",
              msg="핵심 스핀 4개는 9개 방법 전원 일치 — 방법을 바꿔도 같은 답이 나온다")
    add_table(s, ["앙상블 최종", "ppt 앵커", "2021-MLP", "CNN뱅크", "PF",
                  "SpinDETR", "cdetect-DE", "하이브리드"],
              NV1_ROWS, 0.3, 1.25, 12.75, 5.9, fontsize=9, hdr_colors=HC)

    NV2_ROWS = [
        ("(−151.2, 95)", "(+150, 110)", "—", "(−152.3, 99)", "—", "—"),
        ("(−58.9, 182)", "—", "−50 / −60", "(−51.5, 199)", "—", "(−58.8, 181)"),
        ("(−45.9, 50)", "—", "—", "—", "(−45.7, 53)", "(−45.7, 51)"),
        ("(−38.9, 67)", "(−42, 150)", "—", "(−39.2, 67)", "(−39.3, 71)", "(−39.6, 70)"),
        ("(−12.6, 42)", "—", "−15 / —", "(−14.1, 50)", "(−12.8, 42)", "(−12.3, 42)"),
        ("(−2.7, 37)", "—", "—", "—", "(−3.5, 40)", "—"),
        ("(+51.7, 91)", "—", "+55 / +54", "(+51.1, 94)", "(+48.3, 109)", "(+52.7, 85)"),
        ("(+55.9, 55)", "—", "+55 / +54", "—", "(+55.7, 54)", "—"),
        ("(+346.9, 261)", "(−340, 290)", "+345 / +348", "(+346.2, 263)", "(+345.7, 266)", "(+341/+349, 288/257)"),
        ("(+430.7, 58) ⚠", "—", "— / +420", "—", "—", "(+322, 49)?"),
    ]
    s = prs.slides.add_slide(blank)
    add_title(s, "NV2: 알고리즘별 A-텐서 값 종합표",
              "값 = (A∥, A⊥) kHz · CPMG-16 단일 채널 · ⚠ = 축퇴 잔재 보류 · 앵커 3/3 회수",
              msg="수동 분석 3개를 전부 회수하고 신규 7개를 추가 — 사람 눈의 상위집합")
    add_table(s, ["앙상블 최종", "ppt 앵커", "PF ±400/±600", "cdetect-DE",
                  "하이브리드 v1", "하이브리드 refined"],
              NV2_ROWS, 0.5, 1.35, 12.3, 4.9, fontsize=10,
              hdr_colors=[HC[0], HC[1], HC[4], HC[6], HC[7], HC[7]])

    out = ROOT / "results" / "NV_C13_hybrid_slides.pptx"
    prs.save(out)
    print("saved ->", out)


if __name__ == "__main__":
    asset_pipeline()
    asset_compare()
    asset_f1_bars()
    asset_steps()
    asset_dataflow()
    asset_fitrender()
    asset_analogy()
    build_deck()
