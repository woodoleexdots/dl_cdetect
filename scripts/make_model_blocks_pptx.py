"""Native-shape PPT block diagrams of the four models (editable in PowerPoint).

Slide 1 : all four architectures side by side (4 rows) — at-a-glance
Slides 2-5 : one model per slide, larger blocks + explanation bullets

Output: results/NV_C13_model_blocks.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]

C_IN, C_REP, C_NET, C_OUT = "E8E8E8", "DBE9F8", "DFF2DF", "FDEBD0"
ROW_COLORS = {"2021": "8B3A2F", "detr": "38761D", "pf": "1155CC", "hyb": "CC0000"}


def block(slide, x, y, w, h, title, sub=None, fill=C_NET, fs=12, sfs=9):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                 Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    shp.line.color.rgb = RGBColor.from_string("666666")
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.text = title
    p0 = tf.paragraphs[0]
    p0.font.size = Pt(fs)
    p0.font.bold = True
    p0.alignment = PP_ALIGN.CENTER
    p0.font.color.rgb = RGBColor.from_string("1A1A1A")
    if sub:
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(sfs)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor.from_string("7A0000")
    return shp


def arrow(slide, x1, y1, x2, y2, w=2.25):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),
                                      Inches(y1), Inches(x2), Inches(y2))
    conn.line.width = Pt(w)
    conn.line.color.rgb = RGBColor.from_string("404040")
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med",
                                            "len": "med"})
    ln.append(tail)
    return conn


def label(slide, x, y, w, h, title, sub, color, fs=14, sfs=10):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].font.size = Pt(fs)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor.from_string(color)
    p = tf.add_paragraph()
    p.text = sub
    p.font.size = Pt(sfs)
    p.font.color.rgb = RGBColor.from_string("595959")


def chain(slide, y, h, blocks, x0=2.1, gap=0.28, cap_fs=None, cap_h=0.45):
    """blocks: (title, sub, fill, width[, caption]). Draws with arrows and,
    when cap_fs is given, a small explanation textbox under each block."""
    x = x0
    for i, b in enumerate(blocks):
        t, s, f, w = b[:4]
        cap = b[4] if len(b) > 4 else None
        block(slide, x, y, w, h, t, s, f)
        if cap_fs and cap:
            tb = slide.shapes.add_textbox(Inches(x - 0.05), Inches(y + h + 0.02),
                                          Inches(w + 0.1), Inches(cap_h))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Pt(0)
            tf.text = cap
            pgh = tf.paragraphs[0]
            pgh.font.size = Pt(cap_fs)
            pgh.alignment = PP_ALIGN.CENTER
            pgh.font.color.rgb = RGBColor.from_string("444444")
        x_end = x + w
        if i < len(blocks) - 1:
            arrow(slide, x_end + 0.02, y + h / 2, x_end + gap - 0.02, y + h / 2)
        x = x_end + gap


def add_title(slide, text, msg=None):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5),
                                  Inches(0.95))
    tf = tb.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    if msg:
        p = tf.add_paragraph()
        p.text = "▶ " + msg
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string("B45000")


def add_bullets(slide, items, top=5.9, size=13):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.4),
                                  Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + it
        p.font.size = Pt(size)


ROWS = {
    "2021": dict(
        name="2021 뱅크 (Jung et al.)", one="윈도우별 독립 판정 · A 위치만",
        blocks=[("CPMG 신호", "1채널 (700,)", C_IN, 1.5, "실험에서 측정한\n코히어런스 곡선"),
                ("주기 접기", "이미지 13×53", C_REP, 1.6, "후보 주기로 잘라 쌓아\n세로선을 확인"),
                ("MLP ×61", "윈도우별 독립 모델", C_NET, 2.0, "윈도우마다 전담\n분류기가 각자 판정"),
                ("3-class 확률", "(3,)/윈도우", C_NET, 1.7, "스핀 0/1/2개일\n확률을 출력"),
                ("피크 검출", "A 위치만", C_OUT, 1.5, "확률 봉우리만\n후보로 추림")],
        bullets=["작동: 후보 주기마다 신호를 접어 이미지로 만들고, 그 윈도우 전담 MLP가 스핀 유무를 판정 — 61개 모델이 각자 판단",
                 "특장점: 단순하고 빠름 · 2021 논문이 저온 데이터에서 31개 스핀을 찾은 검증된 틀",
                 "한계: 윈도우끼리 정보 공유가 없어 상온 노이즈에 오탐 · B(결합강도)는 별도 회귀 필요"]),
    "detr": dict(
        name="SpinDETR", one="end-to-end 집합 예측 · 물리 힌트 없음 (대조군)",
        blocks=[("신호 3채널", "(3, 700)", C_IN, 1.5, "N=8/16/20 세 측정을\n함께 입력"),
                ("Conv 스템", "(175, 128)", C_NET, 1.5, "신호를 짧은 요약\n조각들로 압축"),
                ("Transformer 인코더 ×4", "(175, 128)", C_NET, 2.3, "조각들 사이의\n패턴 관계를 파악"),
                ("쿼리 10개 + 디코더 ×4", "(10, 128)", C_NET, 2.3, "질문 카드 10장이\n각자 스핀을 탐색"),
                ("스핀 집합", "10×(p, A, B)", C_OUT, 1.6, "카드별 (존재확률,\nA, B)를 출력")],
        bullets=["작동: 질문 카드(쿼리) 10장이 신호 전체에 어텐션해 각자 (존재확률, A, B)를 출력 — Hungarian 매칭으로 학습",
                 "특장점: (A,B) 동시 출력 · 1 forward pass · 고노이즈에서 완만한 성능 저하",
                 "역할: 물리 사전지식 없이 얼마나 되는지 보는 대조군 → PF보다 낮아 사전지식의 가치를 정량화 (−34 kHz 후보의 최강 지지자)"]),
    "pf": dict(
        name="PeriodFormer (제안)", one="물리 사전지식 + 윈도우 간 어텐션 · 오탐 0",
        blocks=[("신호 3채널", "(3, 700)", C_IN, 1.5, "N=8/16/20 세 측정을\n함께 입력"),
                ("61윈도우 접기", "토큰 (61,3,13,53)", C_REP, 1.9, "모든 후보 주기로\n동시에 접음"),
                ("공유 CNN 임베딩", "(61, 128)", C_NET, 1.9, "각 윈도우 그림을\n요약 벡터로 변환"),
                ("어텐션 ×4 (윈도우 간)", "(61, 128)", C_NET, 2.2, "윈도우끼리 증거를\n대조(회의)"),
                ("P(spin) 곡선", "(61,)", C_OUT, 1.6, "A축 위 스핀\n존재 확률 지도")],
        bullets=["작동: 모든 후보 주기로 접은 61개 토큰을 하나의 Transformer가 서로 대조(회의) — 진짜 스핀만 이웃 윈도우에 일관된 흔적",
                 "특장점: 전 벤치마크 오탐 0 · 모델 불일치(자기장 오차 등)에도 precision 1.0 유지",
                 "왜 이렇게: 2021의 실패 원인이 '고립 판정'임을 진단하고, 윈도우들을 토큰으로 만들어 정보를 합치는 지점을 신경망 안으로 옮김"]),
    "hyb": dict(
        name="PF→DE 하이브리드 (최종)", one="신경망 검출 + 물리 피팅 열거 · 전 조건 1위",
        blocks=[("P(spin) 곡선", "PF 출력", C_IN, 1.6, "PF가 만든\n신뢰 지도를 입력"),
                ("문턱 → 후보 영역", "구간 목록", C_REP, 1.9, "확률 높은 구간만\n선별해 탐색 제한"),
                ("영역 제약 DE 피팅", "스핀 +1씩, ∏Mᵢ", C_NET, 2.2, "물리 공식을 데이터에\n맞추며 스핀 추가"),
                ("BIC 선택", "개수 k* 자동", C_NET, 1.6, "추가를 멈출 시점을\n통계 기준이 결정"),
                ("스핀 목록", "{(A∥,A⊥)}×k*", C_OUT, 1.7, "최종 결합강도 목록\nNV1 14개")],
        bullets=["작동: PF가 좁힌 영역 안에서만 물리 공식(∏Mᵢ)을 데이터에 피팅하며 스핀을 하나씩 추가 — BIC가 멈출 지점을 결정",
                 "특장점: 한 영역의 다중 스핀까지 열거 · 실측 50-스핀 검증 3개 조건 모두 1위 (F1 0.78/0.51/0.55)",
                 "왜 이렇게: PF의 약점이 '클러스터 안 개수 세기'임을 진단 → 세는 일은 물리 피팅+통계 기준에 맡기는 분업"]),
}


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- slide 1: all four at a glance ----
    s = prs.slides.add_slide(blank)
    add_title(s, "네 모델 한눈 비교 (블록 다이어그램)",
              msg="회색=입력 · 파랑=표현 · 초록=연산 · 주황=출력 — 세로로 훑으면 설계 차이가 보인다")
    y = 1.35
    for key in ["2021", "detr", "pf", "hyb"]:
        r = ROWS[key]
        label(s, 0.15, y + 0.05, 1.9, 1.1, r["name"], r["one"],
              ROW_COLORS[key], fs=12, sfs=8)
        chain(s, y, 0.85, r["blocks"], cap_fs=7.5, cap_h=0.4)
        y += 1.45
    tb = s.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.4), Inches(0.4))
    tb.text_frame.text = ("하이브리드의 입력 = PeriodFormer의 출력 (3행→4행 연결) · "
                          "괄호 안 붉은 글씨 = 데이터 크기(텐서 shape)")
    tb.text_frame.paragraphs[0].font.size = Pt(11)

    # ---- slides 2-5: one model per slide ----
    for key in ["2021", "detr", "pf", "hyb"]:
        r = ROWS[key]
        s = prs.slides.add_slide(blank)
        add_title(s, r["name"], msg=r["one"])
        n = len(r["blocks"])
        gap = 0.42
        avail = 12.7 - 0.4 - gap * (n - 1)
        scale = avail / sum(b[3] for b in r["blocks"])
        big = [(b[0], b[1], b[2], b[3] * scale) + tuple(b[4:])
               for b in r["blocks"]]
        chain(s, 2.1, 1.6, big, x0=0.4, gap=gap, cap_fs=11, cap_h=0.6)
        add_bullets(s, r["bullets"], top=4.9, size=14)

    out = ROOT / "results" / "NV_C13_model_blocks.pptx"
    prs.save(out)
    print("saved ->", out)


if __name__ == "__main__":
    main()
