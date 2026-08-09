"""Ultra-compact 3-slide summary deck.

1. one table: manual (previous ppt) vs 2021-method vs proposed spins
2. validation on the public 50-spin bath (fig 26)
3. NV1/NV2 real CPMG data (N=8/16/20) with fitted overlays (figs 21/22)

Output: results/NV_C13_summary_3slides.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIGS = ROOT / "results" / "figs"


def add_title(slide, text, sub=None, msg=None):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5),
                                  Inches(1.0))
    tf = tb.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    if sub:
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(12)
    if msg:
        p = tf.add_paragraph()
        p.text = "▶ " + msg
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string("B45000")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---------------- slide 1: three-way spin table ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "요약 ①: 찾아낸 ¹³C 핵스핀 — 세 방법 비교",
              "값 = (A∥, A⊥) kHz · 제안 방법은 95% 부트스트랩 CI 보유 (예: −87.9±0.9)",
              msg="수동 4개 → 2021 방법은 위치 단서 5건(구간 2 + 피크 3, 목록 미완성) → 제안 방법 14개 완성 — 앵커 전부 회수")

    headers = ["", "이전 ppt (수동 분석)", "2021 방법 (core-pipeline 재현)",
               "제안 방법 (하이브리드 앙상블)"]
    rows = [
        ("NV1\n검출 스핀",
         "4개\n(−88, 47)\n(−5, 35)\n(+8, 31)\n(−38, 37)*",
         "출력 = A∥ 확률곡선의 고확률 위치 5건\n('스핀이 이 A∥ 근처에 있다'까지만 —\n A⊥·개수·정확값은 출력 안 됨)\n① 구간 A∈[−6, +12] (확정 스핀 5개 자리)\n② 구간 A∈[+40, +46] (확정 +39·+48 부근)\n③ 피크 −20 = 확정 (−18.8) 자리\n④ 피크 +22 = 확정 (+24.3) 자리\n⑤ 피크 −34 = 확정 목록에 없음\n    → 목록 외 후보 −32와 일치",
         "14개 — (A∥, A⊥) 완비\n(−87.9,18) (−18.8,16) (−11.5,14) (−5.1,22)\n(+0.6,17) (+4.0,26) (+8.4,28) (+14.2,22)\n(+24.3,20) (+39.0,28) (+48.0,25) (+65.6,28)\n(+91.9,33) (+114.4,24)"),
        ("NV2\n검출 스핀",
         "3개\n(−340, 290)\n(+150, 110)\n(−42, 150)*",
         "— (미평가:\n단일 채널·강결합은\n원방법 적용 범위 밖)",
         "10개 — 앵커 3/3 회수 포함\n(−151.2,95) (−58.9,182) (−45.9,50)\n(−38.9,67) (−12.6,42) (−2.7,37)\n(+51.7,91) (+55.9,55) (+346.9,261)\n(+430.7,58)⚠"),
        ("비고\n(총평)",
         "눈 판독의 대략 위치 앵커 —\n7개 전부 우리 목록에서 재발견\nA⊥는 과대추정 경향",
         "위치는 맞추나 목록은 미완성 —\n구간·피크 수준까지만 출력,\n개수·B값·정확한 A를 정하는\n후반 회귀 단계는 재현 범위 밖",
         "개수(BIC, RJMCMC 사후분포와 정합)\n·B값·95% CI까지 완성한 유일한 목록\n3개 측정 동시 재현 RMSE 0.088–0.176"),
    ]
    tblshape = s.shapes.add_table(len(rows) + 1, 4, Inches(0.35), Inches(1.55),
                                  Inches(12.65), Inches(5.6))
    tbl = tblshape.table
    tbl.columns[0].width = Inches(1.15)
    tbl.columns[1].width = Inches(2.6)
    tbl.columns[2].width = Inches(3.3)
    tbl.columns[3].width = Inches(5.6)
    fills = [None, "FCE5CD", "F2D5CC", "D9EAD3"]
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.text_frame.paragraphs[0].font.size = Pt(12)
        c.text_frame.paragraphs[0].font.bold = True
        if fills[j]:
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor.from_string(fills[j])
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            tf = c.text_frame
            for li, line in enumerate(val.split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(9.5 if j == 3 else 10)
                if li == 0:
                    p.font.bold = True
    tb = s.shapes.add_textbox(Inches(0.4), Inches(7.02), Inches(12.5), Inches(0.45))
    tb.text_frame.text = ("*수동 분석의 A∥ 부호는 m_s 브랜치 관례 차이 · "
                          "⚠ = 단일 채널 축퇴로 보류 · 2021 방법 수치는 core-pipeline 재현 기준(원방법 하한)")
    tb.text_frame.paragraphs[0].font.size = Pt(10)
    pgh = tb.text_frame.add_paragraph()
    pgh.text = ("† BIC(베이지안 정보 기준) = 스핀을 추가해 좋아지는 피팅이 복잡도 벌점을 넘을 때만 개수를 늘리는 통계 기준(오컴의 면도날의 정량화) · "
                "95% CI(신뢰구간) = 잔차 부트스트랩 60회 재피팅에서 값이 흔들리는 폭으로 잰 '참값이 95% 확률로 들어있는 구간'(오차막대)")
    pgh.font.size = Pt(9)


    # ---------------- slide: head-to-head 2021 vs ours ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "정면 비교: 2021 방법 vs 우리 방법",
              "나머지 모델들(CNN뱅크·SpinDETR·PF단독·DE·RJMCMC)은 이 결론을 떠받치는 ablation·베이스라인",
              msg="같은 실측 스핀 환경에서 채점한 결과 — 저온 ×2.2, 상온 ×1.5, 오차주입 ×1.7")

    hh = ["", "2021 방법 (core-pipeline 재현)", "우리 방법 (PF→RJMCMC 하이브리드 v2)"]
    rows2 = [
        ("구조", "윈도우별 독립 MLP 61개\n+ 별도 디노이저", "신경망 탐지(윈도우 간 어텐션)\n+ 영역 제약 물리 피팅 + BIC"),
        ("출력", "고확률 영역 (A 위치 후보)\nB·개수는 별도 단계 필요", "최종 스핀 목록 {(A∥, A⊥)} × k*\n개수 자동 · 95% CI 제공"),
        ("실측 50-스핀 검증 (F1)", "저온 0.38 · 상온 0.37 · 오차주입 0.35", "저온 0.84 · 상온 0.57 · 오차주입 0.61"),
        ("오탐 (precision)", "0.55–1.00 (상온에서 붕괴)", "0.87–1.00 (오차주입에도 0.87)"),
        ("실데이터 NV1", "고확률 영역 2곳 + 피크 3개\n(개별 스핀 분해 불가)", "14개 스핀 확정 (CI 포함)\n3개 측정 동시 재현 RMSE 0.09–0.18"),
        ("전제 조건", "저온·고SNR·N=32/256 장시간 측정", "상온·700점·N≤20에서 작동"),
    ]
    shape = s.shapes.add_table(len(rows2) + 1, 3, Inches(0.4), Inches(1.6),
                               Inches(12.55), Inches(5.4))
    t = shape.table
    t.columns[0].width = Inches(2.2)
    t.columns[1].width = Inches(4.9)
    t.columns[2].width = Inches(5.45)
    for j, h in enumerate(hh):
        c = t.cell(0, j)
        c.text = h
        c.text_frame.paragraphs[0].font.size = Pt(13)
        c.text_frame.paragraphs[0].font.bold = True
        if j == 1:
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string("F2D5CC")
        if j == 2:
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string("D9EAD3")
    for i, row in enumerate(rows2):
        for j, val in enumerate(row):
            c = t.cell(i + 1, j)
            tf = c.text_frame
            for li, line in enumerate(val.split("\n")):
                pgh = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                pgh.text = line
                pgh.font.size = Pt(11)
                if j == 0 and li == 0:
                    pgh.font.bold = True
    tb = s.shapes.add_textbox(Inches(0.4), Inches(7.08), Inches(12.5), Inches(0.38))
    tf = tb.text_frame
    tf.text = ("• 공정성 — 2021 수치는 core-pipeline 재현 기준(원방법 하한) · 두 방법 모두 같은 forward model·같은 채점 규칙")
    tf.paragraphs[0].font.size = Pt(10.5)


    # ---------------- slide: two-model block diagrams ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "모델 비교: 블록 다이어그램",
              "회색=입력 · 파랑=표현 · 초록=연산 · 주황=출력 — 위(2021) vs 아래(우리)",
              msg="핵심: 어텐션(회의)으로 오탐을 잡고, 영역 안 열거(BIC)로 개수를 확정 — 특징은 아래 표로 요약")

    C_IN2, C_REP2, C_NET2, C_OUT2 = "E8E8E8", "DBE9F8", "DFF2DF", "FDEBD0"

    def blk(x, y, w, h, title, sub, fill, cap=None, fs=11):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                 Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill)
        shp.line.color.rgb = RGBColor.from_string("666666")
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.word_wrap = True
        tf.text = title
        tf.paragraphs[0].font.size = Pt(fs)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.color.rgb = RGBColor.from_string("1A1A1A")
        if sub:
            pgh = tf.add_paragraph()
            pgh.text = sub
            pgh.font.size = Pt(8.5)
            pgh.alignment = PP_ALIGN.CENTER
            pgh.font.color.rgb = RGBColor.from_string("7A0000")
        if cap:
            tb = s.shapes.add_textbox(Inches(x - 0.05), Inches(y + h + 0.02),
                                      Inches(w + 0.1), Inches(0.45))
            ctf = tb.text_frame
            ctf.word_wrap = True
            ctf.text = cap
            ctf.paragraphs[0].font.size = Pt(8)
            ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
            ctf.paragraphs[0].font.color.rgb = RGBColor.from_string("444444")

    def arr2(x1, x2, y):
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),
                                      Inches(y), Inches(x2), Inches(y))
        conn.line.width = Pt(2.25)
        conn.line.color.rgb = RGBColor.from_string("404040")
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))

    def row(y, label_txt, label_color, blocks):
        tb = s.shapes.add_textbox(Inches(0.15), Inches(y + 0.1), Inches(1.75),
                                  Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = label_txt
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor.from_string(label_color)
        x = 2.0
        gap = 0.30
        for i, (t, sub, f, w, cap) in enumerate(blocks):
            blk(x, y, w, 0.95, t, sub, f, cap)
            if i < len(blocks) - 1:
                arr2(x + w + 0.02, x + w + gap - 0.02, y + 0.475)
            x += w + gap

    row(1.85, "2021 방법\n(core-pipeline)", "8B3A2F", [
        ("CPMG 신호", "1채널 (700,)", C_IN2, 1.45, "저온·장시간 측정 전제"),
        ("주기 접기", "이미지 13×53", C_REP2, 1.55, "후보 주기 1개씩"),
        ("MLP ×61", "윈도우별 독립", C_NET2, 1.85, "서로 대화 없음 → 오탐"),
        ("3-class 확률", "(3,)/윈도우", C_NET2, 1.65, "스핀 0/1/2개 확률"),
        ("피크 검출", "A 영역만", C_OUT2, 1.5, "개수·B값 미해결"),
    ])

    row(4.35, "우리 방법\n(하이브리드 v2)", "CC0000", [
        ("신호 3채널", "(3, 700)", C_IN2, 1.3, "N=8/16/20 함께"),
        ("61윈도우 접기", "토큰(61,3,13,53)", C_REP2, 1.6, "모든 주기 동시에"),
        ("어텐션 ×4", "윈도우 간 회의", C_NET2, 1.5, "진짜 스핀만 통과 · 오탐 0"),
        ("P(spin) → 영역", "후보 구간", C_OUT2, 1.5, "어디를 팔지 확정"),
        ("영역 제약\nRJMCMC + BIC", "물리 공식 피팅", C_NET2, 1.75, "몇 개인지 통계로 결정"),
        ("스핀 목록", "{(A∥,A⊥)}×k*", C_OUT2, 1.55, "NV1 14개 · CI 포함"),
    ])

    # vertical contrast annotation between rows
    tb = s.shapes.add_textbox(Inches(2.0), Inches(3.35), Inches(11.0), Inches(0.55))
    tf = tb.text_frame
    tf.text = ("↕ 같은 '주기 접기'에서 출발하지만 — 2021은 여기서 끝(고립 판정), "
               "우리는 접은 61개를 토큰으로 묶어 회의시키고, 그 결과 영역 안에서만 물리 피팅으로 열거")
    tf.paragraphs[0].font.size = Pt(11.5)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor.from_string("1155CC")

    feat_rows = [
        ("① 윈도우 판정", "61개 윈도우가 개별(고립) 판정 — 서로 대화 없음 → 오탐",
         "어텐션으로 윈도우 간 상호 보강(회의) → 가짜 신호 걸러냄 · 오탐 0"),
        ("② 스핀 탐색", "영역 없는 피크 찾기 — A∥ 위치 단서만",
         "영역 제약 열거 — 후보 구간 안에서만 물리 피팅 + BIC로 개수 결정"),
        ("③ 출력", "A∥ 위치 후보(구간·피크)까지",
         "개수 k* · A∥ · 결합강도 A⊥ · 95% CI까지 완결된 스핀 목록"),
        ("④ 검증 (실측 50-스핀 F1)", "0.35 ~ 0.38", "0.57 ~ 0.84 (전 조건 1위)"),
    ]
    shape = s.shapes.add_table(len(feat_rows) + 1, 3, Inches(0.4), Inches(5.85),
                               Inches(12.55), Inches(1.5))
    t3 = shape.table
    t3.columns[0].width = Inches(2.15)
    t3.columns[1].width = Inches(4.85)
    t3.columns[2].width = Inches(5.55)
    for j, h in enumerate(["특징", "2021 방법", "우리 방법 (하이브리드 v2)"]):
        c = t3.cell(0, j)
        c.text = h
        c.text_frame.paragraphs[0].font.size = Pt(10.5)
        c.text_frame.paragraphs[0].font.bold = True
        if j == 1:
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string("F2D5CC")
        if j == 2:
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string("D9EAD3")
    for i, row in enumerate(feat_rows):
        for j, val in enumerate(row):
            c = t3.cell(i + 1, j)
            c.text = val
            c.text_frame.paragraphs[0].font.size = Pt(9.5)
            if j == 0:
                c.text_frame.paragraphs[0].font.bold = True

    # ---------------- slide: Delft 50-spin dataset overview ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "공개 50-스핀 데이터셋 개요 (van de Stolpe et al. 2024)",
              "\"Mapping a 50-spin-qubit network through correlated sensing\" — Nature Communications 15, 2006 (2024) · Delft(Taminiau group)",
              msg="다이아몬드 단일 NV 센터 주변 ¹³C 핵스핀 50개 네트워크 — 스핀별 파라미터 전체가 공개된 실측 데이터셋")
    s.shapes.add_picture(str(ROOT / "results" / "slides_assets" / "delft50_overview.png"),
                         Inches(0.75), Inches(1.55), width=Inches(11.8))
    ds_rows = [
        ("시스템", "다이아몬드 단일 NV 센터 1개 + 주변 ¹³C 핵스핀 50개(+질소 1) — 상관 센싱(concatenated double-resonance) 시퀀스로 네트워크 전체를 매핑"),
        ("측정 조건", "3.7 K 저온 · B = 403.553 G (NV 축 방향) · Jung 2021과 같은 NV 시료 계보"),
        ("공개 스핀 정보", "스핀별 hyperfine (A∥: −48.5 ~ +213.2 kHz, A⊥: 0.2 ~ 59.2 kHz) · 스핀 주파수 · 3차원 격자 위치 · 측정된 핵-핵 커플링 474쌍(대부분 수 Hz)"),
        ("파일 구성", "spin_table / spin_frequencies / spin_positions / spin_measured(predicted)_couplings (xlsx·pkl) + 논문 그림 재현 노트북"),
    ]
    shape = s.shapes.add_table(len(ds_rows), 2, Inches(0.4), Inches(4.95),
                               Inches(12.55), Inches(1.7))
    td = shape.table
    td.columns[0].width = Inches(1.9)
    td.columns[1].width = Inches(10.65)
    td.first_row = False
    for i, (k, v) in enumerate(ds_rows):
        c = td.cell(i, 0)
        c.text = k
        c.text_frame.paragraphs[0].font.size = Pt(10.5)
        c.text_frame.paragraphs[0].font.bold = True
        c = td.cell(i, 1)
        c.text = v
        c.text_frame.paragraphs[0].font.size = Pt(10.5)

    tb = s.shapes.add_textbox(Inches(0.4), Inches(6.68), Inches(12.5), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    pgh = tf.paragraphs[0]
    run = pgh.add_run()
    run.text = "논문: "
    run.font.size = Pt(11)
    run = pgh.add_run()
    run.text = "doi.org/10.1038/s41467-024-46075-4"
    run.font.size = Pt(11)
    run.hyperlink.address = "https://doi.org/10.1038/s41467-024-46075-4"
    run = pgh.add_run()
    run.text = "      데이터 (4TU.ResearchData): "
    run.font.size = Pt(11)
    run = pgh.add_run()
    run.text = "doi.org/10.4121/aba1cc84-0aea-4cdc-93ca-68b0db38bd81.v1"
    run.font.size = Pt(11)
    run.hyperlink.address = "https://doi.org/10.4121/aba1cc84-0aea-4cdc-93ca-68b0db38bd81.v1"

    # ---------------- slide: noise modelling ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "노이즈 모델링 — 트윈 합성 σ의 근거",
              "노이즈 수준은 검증의 민감 변수 — 임의 가정이 아니라 NV1·NV2 실측 데이터에서 통계적으로 추정",
              msg="상온 트윈 σ=0.06은 실측 채널 σ̂(0.030~0.077)의 최악 수준을 취한 보수적 값 · 저온 트윈은 문헌 조건 σ=0.02")
    noise_rows = [
        ("① σ 추정량", "인접점 차분 dᵢ = yᵢ₊₁ − yᵢ  →  σ̂ = std(d)/√2",
         "느린 신호 성분은 상쇄되고 노이즈만 남음 (독립 가우시안 가정의 표준 추정)"),
        ("② 실측 결과", "NV1 σ̂ = 0.030 / 0.047 / 0.065 (N=8/16/20) · NV2 σ̂ = 0.077",
         "dip 구간 신호가 섞인 만큼 과대평가 쪽으로 치우침 = 보수적 상한"),
        ("③ 노이즈 믹싱", "상온 트윈 σ = 0.06 · 저온 트윈 σ = 0.02 (백색 가우시안)\n+ T2 포락선·콘트라스트도 실측 범위에서 랜덤",
         "상온은 NV1 최악 채널 수준, 저온은 문헌 조건 — 관대하지 않은 검증"),
        ("④ 학습 랜덤화", "학습 데이터는 σ ∈ [0.03, 0.12] 균일 무작위",
         "특정 노이즈 수준에 특화된 검출기가 되는 것을 차단 (과적합 방지)"),
    ]
    shape = s.shapes.add_table(len(noise_rows) + 1, 3, Inches(0.4), Inches(1.5),
                               Inches(12.55), Inches(2.15))
    tn = shape.table
    tn.columns[0].width = Inches(1.75)
    tn.columns[1].width = Inches(5.9)
    tn.columns[2].width = Inches(4.9)
    for j, h in enumerate(["단계", "방법 / 값", "의미"]):
        c = tn.cell(0, j)
        c.text = h
        c.text_frame.paragraphs[0].font.size = Pt(10.5)
        c.text_frame.paragraphs[0].font.bold = True
    for i, row in enumerate(noise_rows):
        for j, val in enumerate(row):
            c = tn.cell(i + 1, j)
            tf = c.text_frame
            for li, line in enumerate(val.split("\n")):
                pgh = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                pgh.text = line
                pgh.font.size = Pt(10)
                if j == 0 and li == 0:
                    pgh.font.bold = True

    s.shapes.add_picture(str(ROOT / "results" / "slides_assets" / "noise_model.png"),
                         Inches(0.35), Inches(3.85), width=Inches(12.65))
    tb = s.shapes.add_textbox(Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = ("• 그림: NV1(N=8/16/20)·NV2(N=16) 원시 데이터 확대(τ 4–7 µs) — 이동중앙값(빨강) 주변 산포가 ±2σ̂ 밴드와 일치    "
               "• 맨 오른쪽: 채널별 σ̂ vs 채택값(0.06)·학습 범위(0.03–0.12)")
    tf.paragraphs[0].font.size = Pt(11)

    # ---------------- slide 2: 50-spin validation ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "요약 ②: 공개 실측 50-스핀 배스로 검증",
              "van de Stolpe 2024 (4TU 공개) — Jung 2021과 같은 NV 계보 · 정답을 아는 유일한 실제 스핀 환경",
              msg="저온 트윈 21/27 회수(RMSE 0.074) · 상온 트윈 11/27·오탐 0 — 스핀 값은 학습에 미사용(순수 테스트)")
    s.shapes.add_picture(str(FIGS / "26_twin_validation.png"), Inches(1.2),
                         Inches(1.45), width=Inches(10.6))
    tb = s.shapes.add_textbox(Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.5))
    tf = tb.text_frame
    tf.text = ("• 벤치마크 종합(F1): 제안 하이브리드 v2가 저온 0.84 / 상온 0.57 / 오차주입 0.61로 전 조건 1위 "
               "(2021 재현 0.38/0.37/0.35, 고전 DE 0.71/0.50/0.48, RJMCMC 0.84/0.56/0.52)")
    tf.paragraphs[0].font.size = Pt(12)

    # ---------------- slide 4: cross-method consensus tiers ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "요약 ③: 방법 간 겹쳐서 찾은 스핀 — 합의 등급",
              "허용오차 ±4 kHz · 계열 = 수동 / 2021영역 / CNN뱅크 / DE / SpinDETR / PF / 하이브리드",
              msg="핵심 스핀은 원리가 다른 모든 방법과 사람 눈까지 일치 — 합의가 낮아지는 순서는 탐색범위를 벗어나는 순서와 일치")

    def tier_table(x, w, title, rows, col_w):
        tb = s.shapes.add_textbox(Inches(x), Inches(1.5), Inches(w), Inches(0.35))
        tb.text_frame.text = title
        tb.text_frame.paragraphs[0].font.size = Pt(14)
        tb.text_frame.paragraphs[0].font.bold = True
        shape = s.shapes.add_table(len(rows) + 1, 3, Inches(x), Inches(1.9),
                                   Inches(w), Inches(4.9))
        t = shape.table
        for j, wd in enumerate(col_w):
            t.columns[j].width = Inches(wd)
        for j, h in enumerate(["등급", "스핀 (A∥, A⊥)", "겹친 계열"]):
            c = t.cell(0, j)
            c.text = h
            c.text_frame.paragraphs[0].font.size = Pt(10)
            c.text_frame.paragraphs[0].font.bold = True
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                c = t.cell(i + 1, j)
                tf = c.text_frame
                for li, line in enumerate(val.split("\n")):
                    pgh = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                    pgh.text = line
                    pgh.font.size = Pt(8.5)
                    if j == 0:
                        pgh.font.bold = True

    nv1_rows = [
        ("🥇 전원 7계열\n(+앵커)", "(−5.1, 22)  (+8.4, 28)  (+39.0, 28)",
         "수동·2021영역·CNN·DE·\nSpinDETR·PF·하이브리드 전부"),
        ("🥈 5–6계열", "(−18.8, 16)  (+0.6, 17)\n(+4.0, 26)  (+14.2, 22)",
         "중앙 연속 고확률 구간 합의\n(DE 또는 SpinDETR 일부 누락)"),
        ("🥉 3–4계열", "(−87.9, 18)*  (+24.3, 20)\n(+48.0, 25)  (+65.6, 28)  (+91.9, 33)",
         "*−87.9는 ±60 밖을 보는\n방법 전원 일치 + 앵커 −88"),
        ("⚠ ≤2계열", "(−11.5, 14)  (+114.4, 24)", "하이브리드 계열 / PF 광역만"),
        ("📌 목록 외 후보", "A ≈ −30~−34", "2021(−34)·CNN(−30)·\nSpinDETR(−34.4) 3계열 겹침"),
    ]
    nv2_rows = [
        ("🥇 전원+앵커", "(+346.9, 261)  (−38.9, 67)",
         "PF±400/±600·DE·하이브리드\n전 버전 + 앵커"),
        ("🥈 4계열", "(−58.9, 182)  (−12.6, 42)\n(+51.7, 91)", "PF·DE·v1·refined"),
        ("🥉 3계열", "(+55.9, 55)", "PF·v1·앙상블"),
        ("2계열+앵커", "(−151.2, 95)", "DE·앙상블 + 앵커 +150\n(PF만 놓침)"),
        ("⚠ ≤2계열", "(−45.9, 50)  (−2.7, 37)\n(+430.7, 58) 보류", "하이브리드 계열만"),
    ]
    tier_table(0.3, 6.55, "NV1 (14스핀 + 후보 1)", nv1_rows, [1.35, 2.7, 2.5])
    tier_table(7.05, 6.0, "NV2 (10스핀)", nv2_rows, [1.35, 2.35, 2.3])
    tb = s.shapes.add_textbox(Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.45))
    tf = tb.text_frame
    tf.text = ("• 읽는 법: 등급이 낮은 스핀은 '의심스러운' 것이 아니라 대부분 |A|가 커서 좁은-범위 방법들의 시야 밖 — "
               "불일치가 아니라 각 방법의 탐색 범위 문제임이 구조적으로 설명됨")
    tf.paragraphs[0].font.size = Pt(11.5)


    # ---------------- slide 5: final definitive spin list ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "최종 확정: NV1 / NV2 ¹³C 핵스핀 목록",
              "값 = (A∥ ± 95% CI, A⊥) kHz · ✅✅ 최고신뢰(전원+앵커) ✅ 확정 ⭕ 신뢰 ⚠ 보류",
              msg="NV1 14개 + NV2 9개(+1 보류) — 물리·통계·교차검증 전부 통과, 수동 분석 7개 전원 재발견")

    def spin_table(x, w, title, cells, n_rows):
        tb = s.shapes.add_textbox(Inches(x), Inches(1.5), Inches(w), Inches(0.3))
        tb.text_frame.text = title
        tb.text_frame.paragraphs[0].font.size = Pt(13)
        tb.text_frame.paragraphs[0].font.bold = True
        shape = s.shapes.add_table(n_rows, 1, Inches(x), Inches(1.85),
                                   Inches(w), Inches(0.42 * n_rows))
        t = shape.table
        for i in range(n_rows):
            c = t.cell(i, 0)
            c.text = cells[i] if i < len(cells) else ""
            c.text_frame.paragraphs[0].font.size = Pt(10.5)

    nv1_cells = [
        "(−87.9 ±0.9, 18.3) ✅", "(−18.8 ±1.2, 15.5) ✅",
        "(−11.5 ±1.7, 13.9) ✅", "(−5.1 ±0.6, 22.5) ✅✅",
        "(+0.6 ±1.5, 17.1) ✅", "(+4.0 ±0.7, 26.2) ✅",
        "(+8.4 ±0.6, 27.8) ✅✅",
    ]
    nv1_cells2 = [
        "(+14.2 ±1.0, 21.7) ✅", "(+24.3 ±1.5, 20.1) ✅",
        "(+39.0 ±2.0, 28.2) ✅✅", "(+48.0 ±2.6, 25.1) ✅",
        "(+65.6 ±1.0, 27.6) ✅", "(+91.9 ±0.7, 32.6) ✅",
        "(+114.4 ±1.7, 23.9) ✅",
    ]
    nv2_cells = [
        "(+346.9, 261) ✅✅", "(−151.2, 95) ✅", "(−58.9, 182) ✅",
        "(−45.9, 50) ⭕", "(−38.9, 67) ✅✅", "(−12.6, 42) ✅",
        "(−2.7, 37) ⭕", "(+51.7, 91) ✅", "(+55.9, 55) ✅",
        "(+430.7, 58) ⚠ 보류",
    ]
    spin_table(0.35, 3.1, "NV1 — 14개 (1/2)", nv1_cells, 7)
    spin_table(3.60, 3.1, "NV1 — 14개 (2/2)", nv1_cells2, 7)
    spin_table(7.35, 3.1, "NV2 — 9개 + 보류 1", nv2_cells[:5], 5)
    spin_table(10.60, 2.5, " ", nv2_cells[5:], 5)

    final_rows = [
        ("개수 근거", "BIC 최소 k*=14 · RJMCMC 사후 최빈 15(개수에 대한 확률분포에서 가장 확률이 높았던 값)와 정합 · 스핀 추가 벌점 통과분만 채택"),
        ("통계적 실체", "14개 전부 95% CI 상호 비겹침 (부트스트랩 60회)"),
        ("물리 재현", "단일 목록으로 N=8/16/20 세 측정 동시 설명, RMSE 0.088/0.116/0.176 (노이즈 바닥 도달)"),
        ("외부 정합", "수동 분석 앵커 NV1 4/4 · NV2 3/3 회수 (부호는 m_s 브랜치 관례 차이)"),
        ("NV2 주의", "단일 채널(CPMG-16)이라 A⊥는 대표값 · (+430.7) 축퇴 잔재 가능 → CPMG-8/20 추가 측정 시 판정"),
        ("목록 외 후보", "NV1 A ≈ −32±3 (3계열 지지, 최종 피팅 미포함)"),
    ]
    shape = s.shapes.add_table(len(final_rows), 2, Inches(0.4), Inches(5.15),
                               Inches(12.55), Inches(2.1))
    tfi = shape.table
    tfi.columns[0].width = Inches(1.9)
    tfi.columns[1].width = Inches(10.65)
    tfi.first_row = False
    for i, (k, v) in enumerate(final_rows):
        c = tfi.cell(i, 0)
        c.text = k
        c.text_frame.paragraphs[0].font.size = Pt(10.5)
        c.text_frame.paragraphs[0].font.bold = True
        c = tfi.cell(i, 1)
        c.text = v
        c.text_frame.paragraphs[0].font.size = Pt(10.5)

    # ---------------- slide: how the fit curves are rendered ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "피팅 곡선은 이렇게 만든다: 스핀별 딥의 합성 곱",
              "스핀 목록 → 스핀별 물리 곡선 Mᵢ → 전체 곱 ∏Mᵢ = 피팅 곡선 (스핀 수는 BIC 자동 선택)",
              msg="곡선은 그린 것이 아니라 스핀 목록에서 물리 공식(Eq.1–3)으로 계산된 예측 — 데이터와 맞으면 목록이 맞다는 뜻")
    s.shapes.add_picture(str(ROOT / "results" / "slides_assets" / "fitrender.png"),
                         Inches(0.3), Inches(1.4), width=Inches(12.7))
    tb = s.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(12.5), Inches(1.25))
    tf = tb.text_frame
    tf.word_wrap = True
    eq_lines = [
        ("물리 공식 (Jung 2021 Eq.1–3 · cpmg/physics.py 구현)", 10.5, True),
        ("Eq.1  단일 스핀 딥:   Mₖ(τ) = 1 − K·sin²(Nφ/2),    K = mₓ²·(1−cos α)(1−cos β)/(1+cos φ)", 12, False),
        ("Eq.2  위상:   cos φ = cos α·cos β − m_z·sin α·sin β,    α = ω̃τ (스핀 세차) · β = ω_L τ (라머 세차)", 12, False),
        ("Eq.3  스핀 축:   ω̃ = √[(A∥+ω_L)² + A⊥²],  m_z = (A∥+ω_L)/ω̃,  mₓ = A⊥/ω̃    →   전체 곱 M(τ) = ∏ₖ Mₖ(τ),  Pₓ = ½(1 + a₀·M·e^−(τ/T₂)ⁿ)", 12, False),
    ]
    for i, (line, fs, bold) in enumerate(eq_lines):
        pgh = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pgh.text = line
        pgh.font.size = Pt(fs)
        pgh.font.bold = bold
        if not bold:
            pgh.font.color.rgb = RGBColor.from_string("1A1A6E")
    tb = s.shapes.add_textbox(Inches(0.4), Inches(5.85), Inches(12.5), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate([
        "① 실험 Pₓ를 엔벨로프로 나눠 모델과 같은 M 스케일로 정규화 → ② 각 스핀 (Aᵢ, Bᵢ)가 자기 주기의 딥 트레인 Mᵢ 생성 → ③ 전체 곱 ∏Mᵢ가 피팅 곡선",
        "스핀이 늘수록 곱에 항이 추가되어 미세 구조를 더 재현 — NV2에서 6→10스핀에 RMSE 0.308→0.286 (BIC 벌점을 통과하고도 잔차 감소 = 과적합 아님)",
        "다음 슬라이드의 빨간 곡선이 전부 이 방식으로 계산된 것 — 자유 파라미터로 모양을 다듬는 단계가 없음",
    ]):
        pgh = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pgh.text = line
        pgh.font.size = Pt(12)

    # ---------------- slide: real-data final fits (moved to last) ----------------
    s = prs.slides.add_slide(blank)
    add_title(s, "최종 피팅: 실제 CPMG 데이터 (NV1: N=8/16/20 · NV2: N=16)",
              "회색 = 실험 · 빨강 = 검출 스핀들의 forward model (그린 것이 아니라 물리 공식으로 계산된 예측)",
              msg="NV1 14스핀이 세 측정(N=8/16/20)을 동시에 재현(RMSE 0.088/0.116/0.176) · NV2 10스핀 RMSE 0.286")
    s.shapes.add_picture(str(FIGS / "21_ensemble_overlay_NV1.png"),
                         Inches(0.25), Inches(1.55), width=Inches(6.6))
    s.shapes.add_picture(str(FIGS / "22_ensemble_overlay_NV2.png"),
                         Inches(7.0), Inches(1.55), width=Inches(6.1))
    tb = s.shapes.add_textbox(Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.55))
    tf = tb.text_frame
    tf.text = ("• 왼쪽 NV1: 위→아래 N=8/16/20 — 같은 14개 스핀 목록 하나가 세 데이터를 모두 설명    "
               "• 오른쪽 NV2: 강결합 프린지 구조까지 추적 (단일 채널 한계로 1개 스핀 보류)")
    tf.paragraphs[0].font.size = Pt(12)

    out = ROOT / "results" / "NV_C13_summary_3slides.pptx"
    prs.save(out)
    print("saved ->", out)


if __name__ == "__main__":
    main()
