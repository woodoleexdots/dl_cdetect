"""Append three slides to the hand-edited 0809 summary decks (kor/eng):

  A. WD method block-by-block, detection stage (blocks 1-4)
  B. WD method block-by-block, confirmation stage (blocks 5-6)
  C. Related work: Poteshman et al. (Quantum 2026, arXiv 2506.18802 / 2506.19259)

Slides A/B are inserted right after the existing block-diagram slide (pos 3),
slide C right after them, preserving all hand-edited content.
"""

from __future__ import annotations

import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PAGE_W, PAGE_H = 13.333, 7.5


def _set_text(tf, lines, size=11, bold_first=True, color=None, align=None):
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.font.size = Pt(size if not (bold_first and i == 0) else size + 1)
        p.font.bold = bold_first and i == 0
        if color is not None:
            p.font.color.rgb = RGBColor.from_string(color)
        if align is not None:
            p.alignment = align


def add_title(s, title, msg):
    tb = s.shapes.add_textbox(Inches(0.35), Inches(0.18), Inches(12.6), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    tb2 = s.shapes.add_textbox(Inches(0.35), Inches(0.78), Inches(12.6), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "▶ " + msg
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor.from_string("B45000")


def add_block(s, x, y, w, h, head, body, fill, head_size=12, body_size=9.5):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    shp.line.color.rgb = RGBColor.from_string("666666")
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = head
    p.font.size = Pt(head_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string("1A1A1A")
    p.alignment = PP_ALIGN.LEFT
    for ln in body:
        pp = tf.add_paragraph()
        pp.text = ln
        pp.font.size = Pt(body_size)
        pp.font.color.rgb = RGBColor.from_string("262626")
    return shp


def add_footer(s, text, y=7.02, size=10):
    tb = s.shapes.add_textbox(Inches(0.35), Inches(y), Inches(12.6), Inches(0.42))
    p = tb.text_frame.paragraphs[0]
    tb.text_frame.word_wrap = True
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor.from_string("404040")


TXT = {
    "kor": dict(
        a_title="WD 방법 상세 ①  — 탐지 단계 (신경망): 블록 1~4",
        a_msg="핵심 아이디어: '주기 가설 61개를 동시에 세우고, 가설끼리 회의시켜 진짜만 남긴다'",
        a_blocks=[
            ("① 입력: 3채널 CPMG 신호  (3×700)",
             ["N=8/16/20 세 측정의 M(τ)를 함께 입력",
              "엔벨로프(Eq.5 신장지수)로 나눠 감쇠 제거 후 M 스케일로 정규화",
              "왜 3채널: 진짜 스핀은 N이 바뀌면 딥 위치·깊이가 물리 법칙대로 변함",
              "→ 우연한 노이즈 딥과 구별하는 결정적 단서"]),
            ("② 61-윈도우 접기 → 토큰  (61,3,13,53)",
             ["후보 A∥ 61개(−60~+60 kHz, 2 kHz 간격)마다",
              "그 A의 목표 주기 T_P(A)로 신호를 잘라 쌓아 13×53 이미지 생성",
              "맞는 주기로 접으면 딥이 세로줄로 정렬, 틀리면 흩어짐 (2021과 동일 표현)",
              "→ '이 A에 스핀이 있나?'라는 가설 61개를 병렬로 세움"]),
            ("③ 공유 CNN 임베딩 + 어텐션 ×4",
             ["공유 CNN: 각 토큰(3×13×53)을 벡터 하나로 압축",
              "— 61개 윈도우가 같은 CNN을 쓰므로 '정렬됨' 판단 기준이 통일됨",
              "어텐션 4층: 윈도우끼리 서로 참조(회의)",
              "— 이웃들이 같이 흥분하면(연속 고확률 구간·하모닉) 문맥으로 진위 판별",
              "→ 2021 독립 MLP의 오탐 원인(서로 대화 없음)을 구조적으로 제거"]),
            ("④ 출력: P(spin) 곡선 → 후보 영역",
             ["윈도우별 스핀 존재 확률 61개 → p>0.35 연속 구간을 margin과 함께 묶음",
              "출력 = 후보 영역 목록 [lo, hi] (예: NV1에서 [−6,+12], [+40,+46] 등)",
              "이 단계의 역할은 '어디를 팔지'까지만 —",
              "개수·정확한 (A∥,A⊥)는 다음 단계(물리 피팅)가 결정"]),
        ],
        a_foot="②는 2021과 같은 표현(주기 접기)을 쓰지만, ③에서 61개를 한 모델이 함께 보는 것이 결정적 차이 — 상온 노이즈에서 오탐 0의 근원",
        b_title="WD 방법 상세 ② — 확정 단계 (물리 피팅): 블록 5~6",
        b_msg="핵심 아이디어: '신경망이 지목한 영역 안에서만, 물리 공식과 통계 기준으로 개수까지 확정한다'",
        b_blocks=[
            ("⑤-1  영역 제약 RJMCMC — 이동 3종",
             ["birth: PF 영역 안에서만 새 스핀 (A,B) 무작위 제안 (20%)",
              "death: 기존 스핀 하나 제거 시도 (20%)",
              "perturb: 개수 유지, A±2 kHz·B±5 kHz 미세 이동 (60%)",
              "영역 밖 제안은 기각 → 탐색 부피 수십 배 축소 = 속도·정확도 동시 확보"]),
            ("⑤-2  채택 규칙 — '월세' 판정",
             ["목적함수: −2lnπ = SSE/σ² + 2k·ln(n)   (n = 3채널×700점)",
              "채택 확률 = min(1, exp(−Δ/2T)) — 좋아지면 무조건, 나빠지면 확률적으로",
              "새 스핀은 월세 2ln(n)≈15.3을 잔차 감소로 내야만 생존 (과적합 차단)",
              "어닐링 T: 20→1 (초반 넓게 탐사 → 후반 엄격 정착)",
              "3채널 동시 잔차 → 한 채널만 맞는 가짜 해는 자동 탈락"]),
            ("⑤-3  개수 결정 — 이중 확인",
             ["BIC 추적: k=1,2,…에서 최소점 k*=14 (NV1)",
              "RJMCMC 사후분포 P(k|데이터): 최빈값 15",
              "서로 독립인 두 절차가 같은 답(14±1) → '개수'의 통계적 근거",
              "출력: MAP 스핀 목록 {(A∥,A⊥)}×k*"]),
            ("⑥  최종 목록 + 95% CI",
             ["잔차 부트스트랩 60회 재피팅 → 각 (A∥,A⊥)의 오차막대(95% CI)",
              "CI가 서로 겹치지 않음 = 통계적으로 구별되는 실체",
              "검증: 앵커 7/7 회수 · NV1 세 측정 동시 재현 RMSE 0.088–0.176",
              "NV1 14개 · NV2 9개+1 보류 확정"]),
        ],
        b_foot="⑤가 하는 일을 한 줄로: '신경망의 직감(영역)'을 '물리학의 심판(forward model 잔차 + 통계 벌점)'으로 확정하는 단계",
        c_title="관련 연구 — Poteshman et al. (Taminiau 그룹 공저, Quantum 2026)",
        c_msg="최신 독립 연구도 같은 결론(trans-dimensional 베이지안이 정답) — 차이는 '탐색을 무엇으로 좁히나': DFT 격자 vs 학습된 어텐션",
        c_rows=[
            ("항목", "Poteshman 2026 (hybrid MCMC)\narXiv 2506.18802 · Quantum 게재",
             "WD 방법 (PF→RJMCMC 하이브리드 v2)"),
            ("탐색 제안\n(사전정보)", "DFT 격자점 이산화 — 다이아몬드 격자 3,518 사이트의\n계산된 하이퍼파인 값 위에서만 탐색 (|A|>5 kHz)",
             "학습된 어텐션 탐지기(PF)의 후보 영역 —\nDFT·격자 정보 불필요, 데이터에서 직접 제안"),
            ("엔진", "RJMCMC + 병렬 템퍼링(10 replica) + RWMH 순환\n(50/100/25 스텝 사이클, 5 앙상블×25,000 스텝)",
             "영역 제약 RJMCMC(birth/death/perturb) + 어닐링\n+ BIC 독립 교차확인"),
            ("감쇠·데이터", "감쇠 λ를 우도에 포함해 동시 추정\n저온 3.7 K · N=32 단일채널 · τ 6–8 µs 250점",
             "신장지수 엔벨로프 사전 정규화 (Eq.5)\n상온 · N=8/16/20 3채널 동시 · 700점"),
            ("결과", "모달 k=46 (참조 50) · 48개 중 45개 posterior 포함\n강결합 1개(C36) 미검출 — DFT 오차 의존\n약결합 <25 kHz는 미해결로 명시",
             "트윈 F1 저온 0.84 / 상온 0.57 / 오차주입 0.61\nNV1 14개 확정(95% CI) · 앵커 7/7 회수\n약결합 kHz 대역이 설계 목표"),
            ("계산 비용", "~75만 forward 평가 · CPU 5코어 ~8시간 · 무학습",
             "PF 학습 1회(GPU 수십 분, 조건별 재사용)\n+ RJMCMC 수 분"),
        ],
        c_foot="자매 논문 arXiv 2506.19259 (Poteshman et al.): 반도체 스핀결함 대량(high-throughput) 특성화 프레임워크 — 우리 자동 파이프라인이 겨냥하는 활용 무대. "
               "두 논문 모두 인용 예정: 방향의 타당성을 독립 확인해 주는 동시대 연구이며, 상온·다채널·학습 제안·완성된 목록(CI)은 우리 쪽 고유 기여.",
    ),
    "eng": dict(
        a_title="WD Method in Detail ① — Detection Stage (Neural): Blocks 1–4",
        a_msg="Key idea: 'stand up 61 period hypotheses at once, let them confer, keep only the real ones'",
        a_blocks=[
            ("① Input: 3-channel CPMG signal  (3×700)",
             ["M(τ) from the N=8/16/20 measurements enters together",
              "Divided by the envelope (Eq.5 stretched-exp) to remove decay",
              "Why 3 channels: a real spin shifts its dips with N following physics",
              "→ the decisive cue separating real spins from noise dips"]),
            ("② 61-window folding → tokens  (61,3,13,53)",
             ["For each of 61 candidate A∥ (−60…+60 kHz, 2 kHz steps),",
              "fold the signal at that A's target period T_P(A) into a 13×53 image",
              "Right period → dips align as vertical stripes; wrong → scattered",
              "(same folding as 2021) → 61 hypotheses stood up in parallel"]),
            ("③ Shared CNN embedding + attention ×4",
             ["Shared CNN: each token (3×13×53) compressed to one vector",
              "— all 61 windows share one CNN, so 'aligned' is judged uniformly",
              "4 attention layers: windows cross-reference each other ('confer')",
              "— joint excitement (continuous high-prob interval / harmonics)",
              "  is resolved by context → removes 2021's isolated-MLP false alarms"]),
            ("④ Output: P(spin) curve → candidate regions",
             ["61 per-window probabilities → contiguous p>0.35 runs merged",
              "Output = candidate regions [lo, hi] (e.g. [−6,+12], [+40,+46] on NV1)",
              "This stage only decides where to dig —",
              "count and exact (A∥,A⊥) are left to the physics-fitting stage"]),
        ],
        a_foot="Block ② uses the same folding as 2021; the decisive difference is ③, where one model sees all 61 windows together — the root of 0 false alarms at room temperature",
        b_title="WD Method in Detail ② — Confirmation Stage (Physics Fitting): Blocks 5–6",
        b_msg="Key idea: 'only inside the network-proposed regions, physics formulas and statistics finalize even the count'",
        b_blocks=[
            ("⑤-1  Region-constrained RJMCMC — 3 moves",
             ["birth: propose a new spin (A,B) only inside PF regions (20%)",
              "death: try removing an existing spin (20%)",
              "perturb: keep count, nudge A±2 kHz · B±5 kHz (60%)",
              "Out-of-region proposals rejected → search volume shrinks",
              "tens-fold = speed and accuracy at once"]),
            ("⑤-2  Acceptance rule — the 'rent' test",
             ["Objective: −2lnπ = SSE/σ² + 2k·ln(n)   (n = 3 channels × 700 pts)",
              "Accept prob = min(1, exp(−Δ/2T)) — always if better, else stochastic",
              "A new spin must pay rent 2ln(n)≈15.3 in residual reduction to survive",
              "Annealing T: 20→1 (broad early exploration → strict settling)",
              "Joint 3-channel residual → single-channel-only fits die out"]),
            ("⑤-3  Deciding the count — double check",
             ["BIC trace over k=1,2,…: minimum at k*=14 (NV1)",
              "RJMCMC posterior P(k|data): mode 15",
              "Two independent procedures, same answer (14±1)",
              "Output: MAP spin list {(A∥,A⊥)}×k*"]),
            ("⑥  Final list + 95% CI",
             ["Residual bootstrap, 60 refits → error bars (95% CI) per (A∥,A⊥)",
              "Non-overlapping CIs = statistically distinct entities",
              "Checks: anchors 7/7 recovered · NV1 three measurements",
              "reproduced simultaneously, RMSE 0.088–0.176",
              "Final: NV1 14 · NV2 9 + 1 withheld"]),
        ],
        b_foot="Block ⑤ in one line: the network's intuition (regions) is finalized by the referee of physics (forward-model residual + statistical penalty)",
        c_title="Related Work — Poteshman et al. (with Taminiau group, Quantum 2026)",
        c_msg="An independent contemporary study reaches the same conclusion (trans-dimensional Bayesian is the right tool) — the difference is what narrows the search: DFT lattice vs learned attention",
        c_rows=[
            ("Aspect", "Poteshman 2026 (hybrid MCMC)\narXiv 2506.18802 · published in Quantum",
             "WD method (PF→RJMCMC hybrid v2)"),
            ("Search proposal\n(prior)", "DFT lattice discretization — search only over computed\nhyperfine values at 3,518 diamond lattice sites (|A|>5 kHz)",
             "Candidate regions from a learned attention detector (PF) —\nno DFT/lattice info needed, proposed from data"),
            ("Engine", "RJMCMC + parallel tempering (10 replicas) + RWMH cycle\n(50/100/25-step cycles, 5 ensembles × 25,000 steps)",
             "Region-constrained RJMCMC (birth/death/perturb) + annealing\n+ independent BIC cross-check"),
            ("Decay · data", "decay λ co-estimated inside the likelihood\ncryo 3.7 K · N=32 single channel · τ 6–8 µs, 250 pts",
             "stretched-exp envelope pre-normalization (Eq.5)\nroom temp · N=8/16/20 three channels jointly · 700 pts"),
            ("Results", "modal k=46 (ref. 50) · 45 of 48 spins in posterior\n1 strong spin (C36) missed — depends on DFT accuracy\nweak coupling <25 kHz stated unresolved",
             "twin F1: cryo 0.84 / room 0.57 / error-injected 0.61\nNV1 14 confirmed (95% CI) · anchors 7/7 recovered\nweak-coupling kHz band is the design target"),
            ("Compute", "~750k forward evaluations · 5 CPU cores, ~8 h · no training",
             "PF trained once (GPU, tens of min; reused per condition)\n+ RJMCMC in minutes"),
        ],
        c_foot="Companion paper arXiv 2506.19259 (Poteshman et al.): a high-throughput spin-defect characterization framework — the arena our automated pipeline targets. "
               "Both will be cited: contemporary work independently validating the direction, while room-temperature operation, multi-channel input, learned proposals, and a completed list with CIs remain our contributions.",
    ),
}

FILLS_A = ["DBE9F8", "DBE9F8", "D9EAD3", "F9E3B3"]
FILLS_B = ["EADCF8", "EADCF8", "F2D5CC", "D9EAD3"]


def four_blocks(s, blocks, fills):
    w, gap, x0, y0, h = 3.08, 0.12, 0.35, 1.30, 5.5
    for i, (head, body) in enumerate(blocks):
        add_block(s, x0 + i * (w + gap), y0, w, h, head, body, fills[i])
        if i < 3:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(x0 + (i + 1) * (w + gap) - gap - 0.01),
                                    Inches(y0 + h / 2 - 0.1),
                                    Inches(gap + 0.02), Inches(0.2))
            ar.fill.solid()
            ar.fill.fore_color.rgb = RGBColor.from_string("888888")
            ar.line.fill.background()


def related_table(s, rows):
    widths = [1.35, 5.6, 5.6]
    tbl = s.shapes.add_table(len(rows), 3, Inches(0.35), Inches(1.30),
                             Inches(sum(widths)), Inches(5.4)).table
    for j, wd in enumerate(widths):
        tbl.columns[j].width = Inches(wd)
    for i, row in enumerate(rows):
        for j, txt in enumerate(row):
            c = tbl.cell(i, j)
            c.text = txt
            for pgh in c.text_frame.paragraphs:
                pgh.font.size = Pt(11 if i == 0 else 9.5)
                pgh.font.bold = (i == 0) or (j == 0)
            if i == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor.from_string(
                    ["EEEEEE", "F2D5CC", "D9EAD3"][j])


def build(lang):
    path = ROOT / "results" / f"NV_C13_summary_0809_{lang}.pptx"
    prs = Presentation(path)
    n0 = len(prs.slides._sldIdLst)
    blank = prs.slide_layouts[6]
    t = TXT[lang]

    s = prs.slides.add_slide(blank)
    add_title(s, t["a_title"], t["a_msg"])
    four_blocks(s, t["a_blocks"], FILLS_A)
    add_footer(s, t["a_foot"])

    s = prs.slides.add_slide(blank)
    add_title(s, t["b_title"], t["b_msg"])
    four_blocks(s, t["b_blocks"], FILLS_B)
    add_footer(s, t["b_foot"])

    s = prs.slides.add_slide(blank)
    add_title(s, t["c_title"], t["c_msg"])
    related_table(s, t["c_rows"])
    add_footer(s, t["c_foot"], y=6.86, size=9.5)

    # move the three appended slides (indices n0..n0+2) to positions 3,4,5
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for off in range(3):
        el = ids[n0 + off]
        lst.remove(el)
        lst.insert(3 + off, el)

    prs.save(path)
    print(f"{lang}: {n0} -> {len(prs.slides._sldIdLst)} slides, saved {path.name}")


if __name__ == "__main__":
    for lang in (sys.argv[1:] or ["kor", "eng"]):
        build(lang)
